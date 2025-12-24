from pptx import Presentation
from pptx.util import Inches

# 创建一个新的PPT演示文稿
prs = Presentation()

# 添加第一张幻灯片
slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 使用标题幻灯片布局
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]

title1.text = "测试PPT标题"
subtitle1.text = "这是第一张幻灯片的副标题"

# 添加第二张幻灯片
slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 使用标题和内容布局
title2 = slide2.shapes.title
title2.text = "第二张幻灯片"

content2 = slide2.shapes.placeholders[1]
text_frame2 = content2.text_frame
text_frame2.text = "这是第二张幻灯片的内容"

# 添加第三张幻灯片
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "第三张幻灯片"

content3 = slide3.shapes.placeholders[1]
text_frame3 = content3.text_frame
text_frame3.text = "这是第三张幻灯片的内容，用于测试PPT转视频功能"

# 保存PPT文件
prs.save("test_ppt.pptx")
print("测试PPT文件已创建：test_ppt.pptx")
