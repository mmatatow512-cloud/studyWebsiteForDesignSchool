import os
import sys
import tempfile
import shutil
import datetime

print("=== PPT转视频完整流程测试 ===")
print(f"当前时间: {datetime.datetime.now()}")
print(f"Python版本: {sys.version}")
print(f"当前工作目录: {os.getcwd()}")

# 添加项目路径到Python路径
sys.path.append(os.path.join(os.getcwd(), 'project'))

# 确保使用正确的Python解释器
print(f"Python解释器路径: {sys.executable}")

# 创建测试PPT
def create_test_ppt(ppt_path, num_slides=3):
    """创建一个简单的测试PPT文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    prs = Presentation()
    
    for i in range(num_slides):
        # 创建幻灯片
        slide_layout = prs.slide_layouts[1]  # 使用带标题和内容的布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = slide.shapes.title
        title.text = f"第{i+1}页 - 测试标题"
        
        # 添加内容
        content = slide.placeholders[1]
        content.text = f"这是第{i+1}页的测试内容。\n\n用于验证PPT转视频功能是否正常工作。"
        
        # 如果有备注页，添加备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            notes.text = f"这是第{i+1}页的备注内容。"
    
    prs.save(ppt_path)
    print(f"✅ 测试PPT创建成功: {ppt_path}")

try:
    # 1. 创建测试PPT
    print("\n--- 1. 创建测试PPT ---")
    temp_dir = tempfile.mkdtemp(prefix="ppt2video_test_")
    print(f"  临时目录: {temp_dir}")
    
    test_ppt = os.path.join(temp_dir, "test_presentation.pptx")
    create_test_ppt(test_ppt, num_slides=2)
    
    # 2. 导入转换器
    print("\n--- 2. 导入转换器 ---")
    from project.ppt2video import ConverterLogic
    print("✅ 转换器导入成功")
    
    # 3. 初始化转换器
    print("\n--- 3. 初始化转换器 ---")
    def log_func(message):
        print(f"[日志] {message}")
    
    converter = ConverterLogic(log_func)
    print("✅ 转换器初始化成功")
    
    # 4. 导出图片
    print("\n--- 4. 导出PPT图片 ---")
    images_folder = os.path.join(temp_dir, "images")
    images = converter.export_images(test_ppt, images_folder)
    if images:
        print(f"✅ 图片导出成功，共 {len(images)} 张图片")
        for img in images:
            size = os.path.getsize(img)
            print(f"  - {os.path.basename(img)}: {size:,} 字节")
    else:
        print("❌ 图片导出失败")
        sys.exit(1)
    
    # 5. 提取文本
    print("\n--- 5. 提取PPT文本 ---")
    scripts = converter.extract_text(test_ppt)
    if scripts:
        print(f"✅ 文本提取成功，共 {len(scripts)} 页")
        for i, script in enumerate(scripts):
            print(f"  - 第{i+1}页: {script[:50]}{'...' if len(script) > 50 else ''}")
    else:
        print("❌ 文本提取失败")
        sys.exit(1)
    
    # 6. 生成音频
    print("\n--- 6. 生成音频 ---")
    audio_folder = os.path.join(temp_dir, "audio")
    audios = converter.generate_audio(scripts, audio_folder, None, None)
    if audios:
        print(f"✅ 音频生成成功，共 {len(audios)} 个音频文件")
        for audio in audios:
            size = os.path.getsize(audio)
            print(f"  - {os.path.basename(audio)}: {size:,} 字节")
    else:
        print("❌ 音频生成失败")
        sys.exit(1)
    
    # 7. 合成视频
    print("\n--- 7. 合成视频 ---")
    output_video = os.path.join(temp_dir, "test_output.mp4")
    success = converter.make_video(images, audios, output_video, use_subtitle=False, scripts=scripts)
    
    if success and os.path.exists(output_video):
        video_size = os.path.getsize(output_video)
        print(f"✅ 视频生成成功！")
        print(f"  文件路径: {output_video}")
        print(f"  文件大小: {video_size:,} 字节 ({video_size/1024:.2f} KB)")
        
        if video_size < 1024:
            print("❌ 警告: 视频文件小于1KB，可能存在问题")
        elif video_size < 10240:
            print("⚠️  注意: 视频文件小于10KB，可能内容较少")
        else:
            print("✅ 视频文件大小正常")
            
        # 复制到桌面方便查看
        desktop_path = os.path.join(os.path.expanduser("~"), "Desktop", "test_ppt2video.mp4")
        shutil.copy2(output_video, desktop_path)
        print(f"📌 测试视频已复制到桌面: {desktop_path}")
        
    else:
        print("❌ 视频生成失败")
        sys.exit(1)
    
except Exception as e:
    print(f"❌ 测试失败: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)
finally:
    # 清理临时文件
    if 'temp_dir' in locals() and os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)
        print(f"\n✅ 临时目录已清理: {temp_dir}")

print(f"\n=== 测试完成 ===")
print(f"测试结束时间: {datetime.datetime.now()}")
