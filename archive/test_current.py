#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
测试当前ppt2video.py的功能是否正常
"""

import os
import sys
import tempfile
import shutil
from pptx import Presentation

sys.path.append(os.path.join(os.path.dirname(__file__), 'project'))

# 导入ppt2video模块
from ppt2video import ConverterLogic

def log_func(message):
    """简单的日志函数"""
    print(f"[LOG] {message}")

def create_test_ppt(output_path, num_slides=3):
    """创建测试PPT文件"""
    prs = Presentation()
    
    for i in range(num_slides):
        # 添加幻灯片
        slide_layout = prs.slide_layouts[1]  # 标题和内容
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = slide.shapes.title
        title.text = f"测试幻灯片 {i+1}"
        
        # 添加内容
        content = slide.placeholders[1]
        content.text = f"这是第 {i+1} 页的内容\n\n用于测试PPT转视频功能"
    
    prs.save(output_path)
    print(f"创建了测试PPT: {output_path}")
    return output_path

def main():
    """主测试函数"""
    print("=== 开始测试当前ppt2video功能 ===")
    
    # 创建临时目录
    temp_dir = tempfile.mkdtemp(prefix="ppt2video_test_")
    print(f"临时目录: {temp_dir}")
    
    try:
        # 创建测试PPT
        test_ppt = os.path.join(temp_dir, "test_presentation.pptx")
        create_test_ppt(test_ppt, num_slides=2)
        
        # 测试ConverterLogic
        converter = ConverterLogic(log_func)
        
        # 步骤1: 导出图片
        print("\n=== 步骤1: 导出PPT图片 ===")
        images_dir = os.path.join(temp_dir, "images")
        os.makedirs(images_dir, exist_ok=True)
        images = converter.export_images(test_ppt, images_dir)
        print(f"导出的图片数量: {len(images)}")
        for img in images:
            print(f"  - {os.path.basename(img)} (大小: {os.path.getsize(img)} 字节)")
        
        # 步骤2: 提取文本
        print("\n=== 步骤2: 提取PPT文本 ===")
        scripts = converter.extract_text(test_ppt)
        print(f"提取的文本数量: {len(scripts)}")
        for i, script in enumerate(scripts):
            print(f"  - 第{i+1}页: {script[:50]}... (长度: {len(script)}字符)")
        
        # 步骤3: 生成音频
        print("\n=== 步骤3: 生成音频文件 ===")
        audio_dir = os.path.join(temp_dir, "audio")
        os.makedirs(audio_dir, exist_ok=True)
        audios = converter.generate_audio(scripts, audio_dir, None, 150)
        print(f"生成的音频数量: {len(audios)}")
        for audio in audios:
            print(f"  - {os.path.basename(audio)} (大小: {os.path.getsize(audio)} 字节)")
        
        # 步骤4: 生成视频
        print("\n=== 步骤4: 生成视频文件 ===")
        output_video = os.path.join(temp_dir, "output_video.mp4")
        success = converter.make_video(images, audios, output_video)
        
        # 验证视频文件
        if success and os.path.exists(output_video):
            video_size = os.path.getsize(output_video)
            print(f"\n=== 测试结果 ===")
            print(f"✅ 视频生成成功！")
            print(f"文件路径: {output_video}")
            print(f"文件大小: {video_size} 字节 ({video_size/1024:.2f} KB)")
            
            if video_size < 1024:
                print("❌ 警告: 视频文件过小，可能存在问题")
            elif video_size < 10240:
                print("⚠️  注意: 视频文件较小，但可能正常")
            else:
                print("✅ 视频文件大小正常")
        else:
            print(f"\n=== 测试结果 ===")
            print(f"❌ 视频生成失败")
            if os.path.exists(output_video):
                print(f"生成的文件大小: {os.path.getsize(output_video)} 字节")
    
    finally:
        # 清理临时文件
        print(f"\n清理临时目录: {temp_dir}")
        shutil.rmtree(temp_dir, ignore_errors=True)

if __name__ == "__main__":
    main()