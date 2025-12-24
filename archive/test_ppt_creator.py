#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建测试用的PPT文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 创建PPT
prs = Presentation()

# 幻灯片1: 标题
slide_layout = prs.slide_layouts[0]  # 标题幻灯片
slide = prs.slides.add_slide(slide_layout)

# 添加标题
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "测试PPT"
subtitle.text = "用于测试PPT转视频功能"

# 幻灯片2: 内容页
slide_layout = prs.slide_layouts[1]  # 标题和内容
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "测试内容页"

content = slide.placeholders[1]
content.text = "这是测试内容\n\n包含多行文本\n用于验证PPT转视频是否正确处理内容"

# 幻灯片3: 图片页
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "图片页"

# 添加文本框
left = Inches(1)
top = Inches(1.5)
width = Inches(6)
height = Inches(1)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "这是一张测试图片的占位符"

# 设置文本格式
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.alignment = PP_ALIGN.CENTER

# 保存PPT
ppt_path = "test.pptx"
prs.save(ppt_path)

print(f"✅ 测试PPT创建成功！")
print(f"文件路径: {os.path.abspath(ppt_path)}")
print(f"幻灯片数量: {len(prs.slides)}")