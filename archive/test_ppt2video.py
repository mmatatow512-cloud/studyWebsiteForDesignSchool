import os
import sys
import tempfile
import traceback

# 添加project目录到Python路径
sys.path.append("c:/Users/23576/Desktop/demo/project")

import ppt2video

def test_ppt_to_video():
    """测试PPT转视频功能的详细流程"""
    print("=== 开始测试PPT转视频功能 ===")
    
    # 定义一个简单的日志函数
    def simple_log(message):
        print(f"[PPT2Video测试] {message}")
    
    try:
        # 1. 准备测试PPT文件
        test_ppt_path = "c:/Users/23576/Desktop/demo/project/test.pptx"
        
        if not os.path.exists(test_ppt_path):
            simple_log("测试PPT文件不存在，创建一个简单的测试PPT")
            # 创建一个简单的PPT文件用于测试
            from pptx import Presentation
            prs = Presentation()
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "测试幻灯片"
            subtitle.text = "这是一个用于测试PPT转视频功能的幻灯片"
            prs.save(test_ppt_path)
            simple_log(f"已创建测试PPT: {test_ppt_path}")
        
        simple_log(f"测试PPT路径: {test_ppt_path}")
        simple_log(f"测试PPT大小: {os.path.getsize(test_ppt_path)} 字节")
        
        # 2. 创建临时目录
        temp_img = tempfile.mkdtemp()
        temp_aud = tempfile.mkdtemp()
        output_video = tempfile.NamedTemporaryFile(suffix='.mp4', delete=False).name
        
        simple_log(f"临时图片目录: {temp_img}")
        simple_log(f"临时音频目录: {temp_aud}")
        simple_log(f"输出视频路径: {output_video}")
        
        # 3. 初始化转换器
        logic = ppt2video.ConverterLogic(simple_log)
        
        # 4. 导出图片
        simple_log("\n第1步：导出PPT幻灯片图片...")
        images = logic.export_images(test_ppt_path, temp_img)
        simple_log(f"导出图片结果: {images}")
        simple_log(f"导出图片数量: {len(images)}")
        
        if images:
            # 检查图片文件
            for i, img_path in enumerate(images):
                if os.path.exists(img_path):
                    simple_log(f"  图片 {i+1}: {img_path}, 大小: {os.path.getsize(img_path)} 字节")
                else:
                    simple_log(f"  图片 {i+1}: {img_path} - 不存在")
        
        # 5. 提取文本
        simple_log("\n第2步：提取PPT文字内容...")
        scripts = logic.extract_text(test_ppt_path)
        simple_log(f"提取文本结果: {scripts}")
        simple_log(f"提取文本数量: {len(scripts)}")
        
        # 6. 生成语音
        simple_log("\n第3步：生成语音文件...")
        voice_id = None
        try:
            import pyttsx3
            engine = pyttsx3.init()
            voices = engine.getProperty('voices')
            simple_log(f"  系统可用语音：{[v.name for v in voices]}")
            for v in voices:
                if "Chinese" in v.name or "CN" in v.name:
                    voice_id = v.id
                    simple_log(f"  选择中文语音：{v.name} (ID: {v.id})")
                    break
            engine.stop()
        except Exception as e:
            simple_log(f"  获取语音失败：{e}")
        
        audios = logic.generate_audio(scripts, temp_aud, voice_id, 170)
        simple_log(f"生成音频结果: {audios}")
        simple_log(f"生成音频数量: {len(audios)}")
        
        if audios:
            # 检查音频文件
            for i, audio_path in enumerate(audios):
                if os.path.exists(audio_path):
                    simple_log(f"  音频 {i+1}: {audio_path}, 大小: {os.path.getsize(audio_path)} 字节")
                else:
                    simple_log(f"  音频 {i+1}: {audio_path} - 不存在")
        
        # 7. 合成视频
        simple_log("\n第4步：合成视频文件...")
        try:
            # 不使用字幕（避免依赖ImageMagick）
            logic.make_video(images, audios, output_video, use_subtitle=False, scripts=None)
            simple_log(f"视频合成完成")
            
            # 检查输出视频文件
            if os.path.exists(output_video):
                file_size = os.path.getsize(output_video)
                simple_log(f"输出视频文件存在: {output_video}")
                simple_log(f"视频文件大小: {file_size} 字节")
                
                if file_size < 10240:
                    simple_log("[错误] 视频文件过小，生成失败")
                else:
                    simple_log("[成功] 视频文件生成成功")
            else:
                simple_log("[错误] 视频文件不存在")
        except Exception as e:
            simple_log(f"[错误] 视频合成失败: {e}")
            traceback.print_exc()
        
    except Exception as e:
        simple_log(f"[严重错误] 测试过程中发生异常: {e}")
        traceback.print_exc()
    
    finally:
        simple_log("\n=== 测试结束 ===")

if __name__ == "__main__":
    test_ppt_to_video()
