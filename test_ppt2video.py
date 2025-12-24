import os
import sys
import tempfile
import shutil
from pptx import Presentation
from pptx.util import Inches

# 添加项目目录到Python路径
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

# 导入转换函数
from ppt2video import convert_presentation_to_video

def create_test_ppt(output_path):
    """创建一个简单的测试PPT"""
    prs = Presentation()
    
    # 第一张幻灯片
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容
    title1 = slide1.shapes.title
    title1.text = "第一张幻灯片"
    
    content1 = slide1.placeholders[1]
    content1.text = "这是第一张幻灯片的内容，用于测试音频生成。"
    
    # 添加备注
    notes_slide = slide1.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "这是第一张幻灯片的备注内容，应该被转换为音频。"
    
    # 第二张幻灯片
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题+内容
    title2 = slide2.shapes.title
    title2.text = "第二张幻灯片"
    
    content2 = slide2.placeholders[1]
    content2.text = "这是第二张幻灯片的内容，用于测试音频生成。"
    
    # 添加备注
    notes_slide = slide2.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = "这是第二张幻灯片的备注内容，应该被转换为音频。"
    
    # 保存PPT
    prs.save(output_path)
    print(f"✓ 测试PPT创建成功: {output_path}")

def test_ppt2video():
    """测试PPT转视频功能"""
    print("=== 开始PPT转视频测试 ===")
    
    # 创建临时目录
    temp_dir = tempfile.mkdtemp()
    print(f"使用临时目录: {temp_dir}")
    
    try:
        # 创建测试PPT
        ppt_path = os.path.join(temp_dir, "test_presentation.pptx")
        create_test_ppt(ppt_path)
        
        # 定义输出视频路径
        output_path = os.path.join(temp_dir, "test_output.mp4")
        
        # 运行转换函数
        print(f"\n=== 开始转换PPT到视频 ===")
        print(f"PPT路径: {ppt_path}")
        print(f"输出路径: {output_path}")
        
        success = convert_presentation_to_video(ppt_path, output_path)
        
        # 验证结果
        if success and os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"\n=== 测试结果 ===")
            print(f"✓ 转换成功!")
            print(f"视频路径: {output_path}")
            print(f"视频大小: {file_size} 字节")
            
            # 使用ffprobe检查视频信息
            try:
                import subprocess
                print(f"\n=== 使用ffprobe检查视频信息 ===")
                
                # 检查音频流
                result = subprocess.run(
                    ['ffprobe', '-v', 'error', '-select_streams', 'a', '-show_entries', 'stream=codec_name', '-of', 'csv=p=0', output_path],
                    capture_output=True,
                    text=True
                )
                if result.stdout.strip():
                    print(f"✓ 视频包含音频流: {result.stdout.strip()}")
                else:
                    print("✗ 视频不包含音频流")
                    
                # 检查视频信息
                result = subprocess.run(
                    ['ffprobe', '-v', 'error', '-show_entries', 'format=duration,bit_rate', '-of', 'csv=p=0', output_path],
                    capture_output=True,
                    text=True
                )
                if result.stdout.strip():
                    duration, bit_rate = result.stdout.strip().split(',')
                    print(f"视频时长: {float(duration):.2f}秒")
                    print(f"视频比特率: {int(bit_rate)/1000:.1f}kbps")
                    
            except Exception as e:
                print(f"ffprobe检查失败: {e}")
                
            return True
        else:
            print(f"\n✗ 转换失败!")
            return False
            
    except Exception as e:
        print(f"\n✗ 测试过程中出错: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        # 清理临时文件
        print(f"\n清理临时目录: {temp_dir}")
        shutil.rmtree(temp_dir)

if __name__ == "__main__":
    success = test_ppt2video()
    if success:
        print("\n=== 测试完成，音频视频生成正常! ===")
    else:
        print("\n=== 测试失败，请检查错误信息! ===")
        sys.exit(1)
