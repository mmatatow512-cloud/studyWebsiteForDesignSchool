#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT转视频集成测试脚本
用于详细测试PPT转视频的完整流程
"""

import os
import sys
import time
import tempfile
import shutil
import traceback
from datetime import datetime

# 添加当前目录到Python路径
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

# 导入需要的模块
import pythoncom
from pptx import Presentation

# 日志函数
def log(message):
    """记录日志"""
    timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_msg = f"[{timestamp}] {message}"
    print(log_msg)
    with open("ppt_integration_test.log", "a", encoding="utf-8") as f:
        f.write(log_msg + "\n")

def create_test_ppt():
    """创建一个简单的测试PPT文件"""
    log("=== 创建测试PPT文件 ===")
    
    test_ppt_path = "test_integration.ppt"
    test_pptx_path = "test_integration.pptx"
    
    # 创建PPTX文件
    try:
        prs = Presentation()
        
        # 添加第一张幻灯片
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片
        slide1.shapes.title.text = "测试幻灯片 1"
        slide1.placeholders[1].text = "这是第一张幻灯片的内容"
        
        # 添加第二张幻灯片
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容
        slide2.shapes.title.text = "测试幻灯片 2"
        slide2.placeholders[1].text = "这是第二张幻灯片的内容\n\n包含多行文本"
        
        # 保存PPTX
        prs.save(test_pptx_path)
        log(f"✓ 创建PPTX测试文件成功: {test_pptx_path}")
        
        # 使用win32com.client创建PPT文件（如果需要）
        # 注意：这个操作可能需要安装Office
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(os.path.abspath(test_pptx_path))
            presentation.SaveAs(os.path.abspath(test_ppt_path), 1)  # 1 = PPT格式
            presentation.Close()
            # powerpoint.Quit()
            log(f"✓ 创建PPT测试文件成功: {test_ppt_path}")
            
            return test_pptx_path  # 返回PPTX路径，兼容性更好
            
        except Exception as e:
            log(f"✗ 创建PPT文件失败 (可能需要安装Office): {e}")
            log("  将使用PPTX格式进行测试")
            return test_pptx_path
            
    except Exception as e:
        log(f"✗ 创建测试PPT失败: {e}")
        traceback.print_exc()
        return None

def test_ppt_conversion(ppt_path):
    """测试PPT转视频的完整流程"""
    log(f"\n=== 测试PPT转视频流程: {ppt_path} ===")
    
    # 导入convert_presentation_to_video函数
    from ppt2video import convert_presentation_to_video
    
    # 设置输出路径
    output_path = "ppt_integration_output.mp4"
    
    # 确保输出路径不存在
    if os.path.exists(output_path):
        os.remove(output_path)
        log(f"已删除旧的输出文件: {output_path}")
    
    # 测试转换
    log(f"开始转换PPT到视频...")
    log(f"PPT路径: {ppt_path}")
    log(f"输出路径: {output_path}")
    
    start_time = time.time()
    try:
        result = convert_presentation_to_video(ppt_path, output_path)
        end_time = time.time()
        
        log(f"转换结果: {result}")
        log(f"转换耗时: {end_time - start_time:.2f} 秒")
        
        # 验证输出文件
        if os.path.exists(output_path):
            size = os.path.getsize(output_path)
            log(f"✓ 输出文件存在: {output_path}")
            log(f"  文件大小: {size} 字节 ({size/1024:.2f} KB)")
            
            if size < 1024:
                log(f"⚠️  文件大小异常，只有 {size} 字节")
                
                # 尝试读取文件内容
                try:
                    with open(output_path, "rb") as f:
                        content = f.read()
                    log(f"  文件内容: {repr(content)}")
                except Exception as e:
                    log(f"  读取文件内容失败: {e}")
            else:
                log(f"✓ 文件大小正常")
        else:
            log(f"✗ 输出文件不存在")
            
    except Exception as e:
        log(f"✗ 转换过程中发生异常: {e}")
        traceback.print_exc()
        
    # 返回结果
    return os.path.exists(output_path) and os.path.getsize(output_path) > 1024

def main():
    """主函数"""
    log("=== PPT转视频集成测试开始 ===")
    
    # 初始化COM组件
    pythoncom.CoInitialize()
    
    # 第一步：创建测试PPT
    ppt_path = create_test_ppt()
    if not ppt_path:
        log("无法创建测试PPT，退出测试")
        return
    
    # 第二步：测试转换
    success = test_ppt_conversion(ppt_path)
    
    log("\n=== PPT转视频集成测试完成 ===")
    log(f"测试结果: {'成功' if success else '失败'}")
    log("详细日志请查看: ppt_integration_test.log")

if __name__ == "__main__":
    main()
