import os
import sys
import tempfile
import shutil

# 设置Python路径
sys.path.append(os.path.abspath('.'))

# 使用正确的Python路径
python_path = sys.executable

print("=== PPT转视频最终测试 ===")
print(f"使用Python: {python_path}")
print(f"当前工作目录: {os.getcwd()}")

# 创建测试目录
test_dir = tempfile.mkdtemp(prefix="ppt2video_final_")
print(f"创建测试目录: {test_dir}")

# 创建一个简单的Python脚本用于测试
script_content = f'''
import os
import sys
import tempfile

# 导入核心函数 - 使用项目的绝对路径确保能找到模块
sys.path.append(r"{os.path.abspath('.')}")
from ppt2video import convert_presentation_to_video

# 创建测试PPT
def create_test_ppt(output_path):
    from pptx import Presentation
    prs = Presentation()
    
    # 幻灯片1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "最终测试"
    subtitle.text = "PPT转视频功能测试"
    
    # 幻灯片2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    content = slide2.placeholders[1]
    title2.text = "测试内容"
    content.text = "这是测试幻灯片内容"
    
    prs.save(output_path)
    return True

# 主测试
if __name__ == "__main__":
    print("=== PPT转视频最终测试 ===")
    
    # 创建测试PPT
    ppt_path = "test_ppt.pptx"
    if not os.path.exists(ppt_path):
        print(f"创建测试PPT: {ppt_path}")
        if not create_test_ppt(ppt_path):
            print("❌ 创建PPT失败")
            sys.exit(1)
    
    print(f"测试PPT: {ppt_path} (大小: {os.path.getsize(ppt_path)} 字节)")
    
    # 输出路径
    output_path = "final_output.mp4"
    print(f"输出视频: {output_path}")
    
    # 执行转换
    print("开始转换...")
    import time
    start_time = time.time()
    
    success = convert_presentation_to_video(ppt_path, output_path)
    
    end_time = time.time()
    elapsed = end_time - start_time
    
    print(f"转换完成，耗时: {elapsed:.2f} 秒")
    print(f"转换结果: {'成功' if success else '失败'}")
    
    # 检查输出
    if os.path.exists(output_path):
        size = os.path.getsize(output_path)
        print(f"✅ 视频文件创建成功")
        print(f"文件大小: {size} 字节 ({size/1024:.2f} KB)")
        if size > 10240:  # 大于10KB视为正常
            print("✅ 文件大小正常")
        else:
            print("⚠️ 文件大小异常")
    else:
        print("❌ 视频文件未创建")
'''

# 写入测试脚本
test_script = os.path.join(test_dir, "final_test.py")
with open(test_script, 'w', encoding='utf-8') as f:
    f.write(script_content)

print(f"创建测试脚本: {test_script}")

# 执行测试脚本
try:
    # 切换到测试目录
    os.chdir(test_dir)
    print(f"切换到测试目录: {os.getcwd()}")
    
    # 执行测试
    print("\n=== 执行测试 ===")
    import subprocess
    result = subprocess.run([python_path, "final_test.py"], capture_output=True, text=True, encoding='utf-8')
    
    # 输出结果
    print("\n=== 测试输出 ===")
    print(result.stdout)
    
    if result.stderr:
        print("\n=== 错误输出 ===")
        print(result.stderr)
    
    print(f"\n=== 测试结束 ===")
    print(f"返回码: {result.returncode}")
    
    # 检查是否生成了视频文件
    output_video = os.path.join(test_dir, "final_output.mp4")
    if os.path.exists(output_video):
        size = os.path.getsize(output_video)
        print(f"\n✅ 测试视频文件: {output_video}")
        print(f"文件大小: {size} 字节 ({size/1024:.2f} KB)")
        if size > 10240:  # 大于10KB视为正常
            print("🎉 测试成功！PPT转视频功能正常工作！")
        else:
            print("⚠️ 测试失败：视频文件仍然太小")
    else:
        print("❌ 测试失败：视频文件未生成")
        
finally:
    # 清理
    print(f"\n清理测试目录: {test_dir}")
    os.chdir(os.path.abspath('.'))
    # shutil.rmtree(test_dir)  # 注释掉，以便手动检查测试结果
    
    print("\n=== 测试完成 ===")
