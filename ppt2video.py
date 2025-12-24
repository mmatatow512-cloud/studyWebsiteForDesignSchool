import pythoncom  # <--- 必须加这个
import tkinter as tk
from tkinter import ttk, filedialog, scrolledtext, messagebox
import threading
import os
import sys
import json
import shutil
import win32com.client
from pptx import Presentation
# 直接导入所需的组件，避免使用不存在的moviepy.editor模块
from moviepy.video.io.VideoFileClip import VideoFileClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
from moviepy.video.VideoClip import ImageClip
from moviepy.video.compositing.CompositeVideoClip import CompositeVideoClip, concatenate_videoclips
from moviepy.video.VideoClip import TextClip, ColorClip
from moviepy.audio.AudioClip import AudioClip
from moviepy.audio.AudioClip import AudioArrayClip
# 如果你需要字幕且安装了ImageMagick，请取消下一行的注释并配置路径

# ===========================
# 配置区域 (防止 ImageMagick 找不到)
# ===========================
# 如果你安装了 ImageMagick，请在这里填入真实路径
magick_path = r"C:\Program Files\ImageMagick-7.1.1-Q16-HDRI\magick.exe"
if os.path.exists(magick_path):
    os.environ['IMAGEMAGICK_BINARY'] = magick_path

# 确保FFmpeg可用
import moviepy.config as mp_config
try:
    # 检查FFmpeg配置
    if hasattr(mp_config, 'FFMPEG_BINARY'):
        print(f"FFMPEG_BINARY配置: {mp_config.FFMPEG_BINARY}")
        # 如果没有配置，尝试自动查找
        if not mp_config.FFMPEG_BINARY or not os.path.exists(mp_config.FFMPEG_BINARY):
            print("FFmpeg未找到，尝试自动配置...")
            # 尝试在PATH中查找FFmpeg
            import subprocess
            try:
                ffmpeg_path = subprocess.check_output(['where', 'ffmpeg'], shell=True, text=True).strip()
                print(f"自动找到FFmpeg: {ffmpeg_path}")
                os.environ['FFMPEG_BINARY'] = ffmpeg_path
                mp_config.FFMPEG_BINARY = ffmpeg_path
            except subprocess.CalledProcessError:
                print("警告: 无法在PATH中找到FFmpeg，请确保已安装并添加到环境变量")
except Exception as e:
    print(f"FFmpeg配置检查失败: {e}")


# ===========================
# 核心逻辑类
# ===========================
class ConverterLogic:
    def __init__(self, logger_func):
        self.log = logger_func  # 用于向界面输出日志

    def export_images(self, ppt_path, temp_folder):
        # 导入必要的模块
        import os
        if not os.path.exists(temp_folder):
            os.makedirs(temp_folder)

        self.log(f"正在启动 PowerPoint 导出图片...")
        self.log(f"PPT文件路径: {ppt_path}")
        self.log(f"PPT文件是否存在: {os.path.exists(ppt_path)}")
        if os.path.exists(ppt_path):
            self.log(f"PPT文件大小: {os.path.getsize(ppt_path)} 字节")
        
        abs_ppt_path = os.path.abspath(ppt_path)
        abs_output_folder = os.path.abspath(temp_folder)

        try:
            self.log(f"创建PowerPoint COM对象...")
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            self.log(f"PowerPoint COM对象创建成功")
            
            # 设置PowerPoint可见性，方便调试
            powerpoint.Visible = 1
            
            self.log(f"打开PPT文件: {abs_ppt_path}")
            # 增加文件存在性和权限检查
            if not os.path.exists(abs_ppt_path):
                self.log(f"[错误] PPT文件不存在: {abs_ppt_path}")
                return []
            
            if not os.access(abs_ppt_path, os.R_OK):
                self.log(f"[错误] 无读取权限: {abs_ppt_path}")
                return []
            
            # 打开PPT文件，增加超时和错误处理
            try:
                self.log(f"[DEBUG] 准备调用PowerPoint.Open: {abs_ppt_path}")
                # 先尝试检查文件的二进制签名，确保是有效的PPT文件
                with open(abs_ppt_path, 'rb') as f:
                    header = f.read(8)
                    self.log(f"[DEBUG] 文件头8字节: {header.hex()}")
                    # PPTX文件以PK开头 (zip文件)
                    if header.startswith(b'PK'):
                        self.log(f"[DEBUG] 确认是PPTX文件 (ZIP格式)")
                    # PPT文件的魔术数字
                    elif header.startswith(b'\xd0\xcf\x11\xe0'):
                        self.log(f"[DEBUG] 确认是PPT文件 (OLE格式)")
                    else:
                        self.log(f"[警告] 可能不是标准PPT/PPTX文件")
                
                presentation = powerpoint.Presentations.Open(abs_ppt_path, WithWindow=False)
                self.log(f"PPT文件打开成功")
            except Exception as e:
                self.log(f"[错误] 打开PPT文件失败: {e}")
                import traceback
                traceback.print_exc()
                return []
            
            # 检查演示文稿属性
            self.log(f"PPT文件名: {presentation.Name}")
            self.log(f"PPT完整路径: {presentation.FullName}")
            self.log(f"PPT标题: {presentation.Title if hasattr(presentation, 'Title') else 'N/A'}")
            self.log(f"PPT创建时间: {presentation.CreationDate if hasattr(presentation, 'CreationDate') else 'N/A'}")
            self.log(f"PPT最后修改时间: {presentation.LastSaveTime if hasattr(presentation, 'LastSaveTime') else 'N/A'}")
            
            # 检查幻灯片集合
            try:
                slides_count = len(presentation.Slides)
                self.log(f"PPT幻灯片数量: {slides_count}")
                
                if slides_count == 0:
                    self.log(f"[警告] PPT中没有幻灯片")
                    return []
            except Exception as e:
                self.log(f"[错误] 获取幻灯片数量失败: {e}")
                return []

            # 导出幻灯片图片
            exported_images = []
            for i, slide in enumerate(presentation.Slides):
                img_path = os.path.join(abs_output_folder, f"slide_{i + 1:02d}.jpg")
                self.log(f"  - 导出第 {i + 1} 页幻灯片到: {img_path}")
                
                try:
                    # 导出高清图 (高度 1080)
                    slide.Export(img_path, "JPG", 0, 1080)
                    
                    # 验证图片是否成功导出
                    if os.path.exists(img_path):
                        img_size = os.path.getsize(img_path)
                        if img_size > 1024:  # 大于1KB才视为有效
                            self.log(f"  - 第 {i + 1} 页幻灯片导出成功，文件大小: {img_size} 字节")
                            exported_images.append(img_path)
                        else:
                            self.log(f"  - [警告] 第 {i + 1} 页幻灯片导出文件太小: {img_size} 字节")
                    else:
                        self.log(f"  - [错误] 第 {i + 1} 页幻灯片导出失败，文件不存在")
                except Exception as e:
                    self.log(f"  - [错误] 导出第 {i + 1} 页幻灯片失败: {e}")
                    continue
            
            self.log(f"共尝试导出 {len(presentation.Slides)} 张幻灯片")
            self.log(f"成功导出并验证 {len(exported_images)} 张有效幻灯片图片")

            presentation.Close()
            self.log(f"PPT文件关闭成功")
            
            # 获取导出的图片列表
            exported_images = sorted([os.path.join(temp_folder, f) for f in os.listdir(temp_folder) if f.endswith(".jpg")])
            self.log(f"共导出 {len(exported_images)} 张幻灯片图片")
            
            # 验证每张图片的大小
            valid_images = []
            for img in exported_images:
                if os.path.exists(img):
                    img_size = os.path.getsize(img)
                    if img_size > 1024:  # 大于1KB才视为有效
                        valid_images.append(img)
                        self.log(f"  - 有效图片: {os.path.basename(img)}，大小: {img_size} 字节")
                    else:
                        self.log(f"  - [警告] 图片太小: {os.path.basename(img)}，大小: {img_size} 字节")
                else:
                    self.log(f"  - [错误] 图片不存在: {os.path.basename(img)}")
            
            self.log(f"有效图片数量: {len(valid_images)}")
            
            # 如果没有有效图片，明确返回错误
            if not valid_images:
                self.log("[错误] 没有成功导出任何有效幻灯片图片")
                return []
                
            return valid_images
        except Exception as e:
            self.log(f"[错误] PPT导出失败: {e}")
            import traceback
            traceback.print_exc()
            return []

    def extract_text(self, ppt_path):
        import os
        import re
        
        # 获取文件扩展名
        file_ext = os.path.splitext(ppt_path)[1].lower()
        
        # 检查文件格式
        if file_ext == '.ppt':
            # 对于旧的.ppt格式，无法使用pptx库提取文字
            self.log("警告: 检测到旧的.ppt格式文件，无法提取文字内容，将使用默认文本")
            # 返回空列表，让调用者处理
            return []
        elif file_ext != '.pptx':
            self.log(f"警告: 不支持的文件格式: {file_ext}，无法提取文字内容")
            return []
            
        self.log("正在解析 PPT 文字内容...")
        try:
            prs = Presentation(ppt_path)
            scripts = []
            for i, slide in enumerate(prs.slides):
                text_content = ""
                # 优先读备注
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    text_content = slide.notes_slide.notes_text_frame.text
                else:
                    # 读正文
                    texts = []
                    if slide.shapes.title:
                        texts.append(slide.shapes.title.text)
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            texts.append(shape.text_frame.text)
                    text_content = "，".join(texts)

                # 增强文本清理：移除重复标点、特殊字符等
                text_content = text_content.replace("\n", "，").replace("\r", "")
                # 移除所有特殊符号
                text_content = re.sub(r'[•●■◆★☆▲▼△▽◁▷◀▶♠♥♦♣]+', '', text_content)
                # 统一替换所有中文标点为逗号
                text_content = re.sub(r'[，。、；：“”‘’"\'！？]+', '，', text_content)
                # 移除重复的逗号
                text_content = re.sub(r'[，]+', '，', text_content)
                # 移除首尾的逗号
                text_content = re.sub(r'^，+|，+$', '', text_content)
                # 移除多余的空格
                text_content = re.sub(r'\s+', '', text_content)
                # 限制长度
                if len(text_content) > 100:
                    text_content = text_content[:97] + "..."
                if not text_content.strip():
                    text_content = "空白幻灯片"  # 空音频占位

                scripts.append(text_content)
            return scripts
        except Exception as e:
            self.log(f"提取文字失败: {e}")
            return []

    def generate_audio(self, scripts, temp_folder, voice_id, rate):
        # 导入必要的模块
        import os
        import sys
        import json
        import tempfile
        import subprocess
        import re
        import threading
        if not os.path.exists(temp_folder):
            os.makedirs(temp_folder)

        audio_paths = []
        self.log("正在生成音频文件...")
        
        # 确保gTTS已安装
        self.log("[调试] 检查并确保gTTS已安装...")
        try:
            from gtts import gTTS
            self.log("[调试] ✓ gTTS已安装")
        except ImportError:
            self.log("[调试] ✗ gTTS未安装，正在安装...")
            try:
                import subprocess
                # 使用python -m pip安装gTTS
                result = subprocess.run(
                    [sys.executable, '-m', 'pip', 'install', 'gtts'],
                    capture_output=True,
                    text=True,
                    timeout=60
                )
                if result.returncode == 0:
                    self.log("[调试] ✓ gTTS安装成功")
                else:
                    self.log(f"[调试] ✗ gTTS安装失败: {result.stderr[:200]}...")
                    # 继续尝试，可能安装过程中有警告但实际成功了
                    try:
                        from gtts import gTTS
                        self.log("[调试] ✓ 尽管有警告，gTTS仍可使用")
                    except ImportError:
                        self.log("[调试] ✗ gTTS确实无法安装")
            except Exception as e:
                self.log(f"[调试] ✗ 安装gTTS时发生异常: {e}")
        
        # 导入所需库
        import numpy as np
        import pyttsx3
        import threading
        import queue
        import re
        import subprocess
        import time
        import tempfile
        import json
        import sys
        
        self.log(f"    设置语速: {rate or 150}")
        self.log(f"    使用语音生成系统: pyttsx3 + 超时保护")
        
        # 创建临时脚本内容模板 (定义在函数顶部，解决作用域问题)
        script_template = '''
import pyttsx3
import os
import sys
import json
import threading
import re
import time

# 获取参数
text = {text}
filename = {filename}
rate = {rate}

result = {result_init}

# 详细的语音信息记录
voice_info = []

# 进一步清理文本，移除过多的标点和省略号
cleaned_text = text
# 移除重复的标点符号
cleaned_text = re.sub(r'[，。、；："\'\n\r]+', '，', cleaned_text)
# 移除重复的顿号
cleaned_text = re.sub(r'[，]+', '，', cleaned_text)
# 移除省略号
cleaned_text = re.sub(r'[.]+', '', cleaned_text)
# 移除特殊符号
cleaned_text = re.sub(r'[•●■◆]', '', cleaned_text)
# 限制文本长度
if len(cleaned_text) > 100:
    cleaned_text = cleaned_text[:100] + "..."

if not cleaned_text.strip():
    cleaned_text = "空白幻灯片"

# 子进程内部的超时处理
class TimeoutError(Exception):
    pass

# 超时处理函数
timeout_occurred = False
def timeout_handler():
    global timeout_occurred
    timeout_occurred = True
    raise TimeoutError("音频生成超时")

def audio_generation_thread():
    """在单独的线程中生成音频"""
    global result
    
    try:
        # 初始化引擎
        engine = pyttsx3.init()
        engine.setProperty('rate', rate)
        
        # 设置中文语音 - 增强的中文语音检测逻辑
        voices = engine.getProperty('voices')
        
        # 记录所有可用语音信息
        for voice in voices:
            voice_info.append({
                "id": voice.id,
                "name": voice.name,
                "languages": voice.languages if hasattr(voice, 'languages') else [],
                "gender": voice.gender if hasattr(voice, 'gender') else "unknown",
                "age": voice.age if hasattr(voice, 'age') else "unknown"
            })
        
        # 增强的中文语音检测
        chinese_voice_found = False
        best_voice = None
        
        # 优先匹配中文语音
        for voice in voices:
            voice_id_lower = voice.id.lower()
            voice_name_lower = voice.name.lower()
            
            # 中文语音关键词匹配
            if any(keyword in voice_id_lower or keyword in voice_name_lower 
                   for keyword in ['chinese', 'zh', 'mandarin', '中文', '普通话', 'cantonese', 'yue']):
                best_voice = voice
                chinese_voice_found = True
                break
        
        if not chinese_voice_found and voices:
            # 如果没有找到中文语音，使用第一个可用语音
            best_voice = voices[0]
            result["warning"] = "未找到中文语音包，使用默认语音: " + best_voice.name
        
        if best_voice:
            engine.setProperty('voice', best_voice.id)
            result["voice_used"] = best_voice.name
        
        # 生成音频
        engine.save_to_file(cleaned_text, filename)
        engine.runAndWait()
        
        # 验证文件
        if os.path.exists(filename):
            file_size = os.path.getsize(filename)
            if file_size > 1024:
                result["success"] = True
                result["message"] = "生成成功，大小: " + str(file_size) + " 字节"
                result["voice_info"] = voice_info
            else:
                result["success"] = False
                result["message"] = "文件太小 (" + str(file_size) + " 字节)"
        else:
            result["success"] = False
            result["message"] = "文件未生成"
            
    except Exception as e:
        result["success"] = False
        result["message"] = "生成错误: " + str(e)
        result["voice_info"] = voice_info

try:
    
    # 使用threading.Timer实现超时
    timeout_timer = threading.Timer(10.0, timeout_handler)
    timeout_timer.start()
    
    try:
        # 启动音频生成线程
        audio_thread = threading.Thread(target=audio_generation_thread)
        audio_thread.start()
        
        # 等待音频生成完成
        audio_thread.join(10.0)
        
        if audio_thread.is_alive():
            # 线程仍在运行，说明超时
            result["success"] = False
            result["message"] = "子进程生成超时"
            # 无法强制终止线程，只能返回错误
    except TimeoutError:
        result["success"] = False
        result["message"] = "子进程生成超时"
    finally:
        # 取消超时计时器
        timeout_timer.cancel()
    
except Exception as e:
    result["success"] = False
    result["message"] = "生成错误: " + str(e)
    result["voice_info"] = voice_info
    # 确保引擎被正确关闭
    try:
        if 'engine' in locals():
            engine.stop()
    except:
        pass

# 输出结果
print(json.dumps(result, ensure_ascii=False))
'''
        
        def generate_audio_with_gtts(text, filename):
            """使用gTTS生成中文语音"""
            try:
                from gtts import gTTS
                import os
                import tempfile
                
                # 清理文本
                cleaned_text = re.sub(r'[，。、；："\'\n\r]+', '，', text)
                cleaned_text = re.sub(r'[•●■◆]', '', cleaned_text)
                cleaned_text = re.sub(r'[，]+', '，', cleaned_text)
                cleaned_text = re.sub(r'^[，]+|[，]+$', '', cleaned_text)
                if not cleaned_text.strip():
                    cleaned_text = "空白幻灯片"
                if len(cleaned_text) > 100:
                    cleaned_text = cleaned_text[:97] + "..."
                
                # 使用gTTS生成语音
                tts = gTTS(text=cleaned_text, lang='zh-cn', slow=False)
                
                # 先保存到临时文件，再移动到目标位置
                temp_file = tempfile.mktemp(suffix='.mp3')
                tts.save(temp_file)
                
                # 转换为wav格式
                if not filename.endswith('.wav'):
                    wav_filename = os.path.splitext(filename)[0] + '.wav'
                else:
                    wav_filename = filename
                
                # 使用ffmpeg转换格式（如果可用）
                try:
                    import subprocess
                    subprocess.run(
                        ['ffmpeg', '-i', temp_file, '-acodec', 'pcm_s16le', '-ar', '44100', wav_filename],
                        check=True,
                        stderr=subprocess.PIPE,
                        stdout=subprocess.PIPE
                    )
                    os.unlink(temp_file)
                except Exception as e:
                    # 如果ffmpeg不可用，使用moviepy转换
                    try:
                        from moviepy.editor import AudioFileClip
                        audio = AudioFileClip(temp_file)
                        audio.write_audiofile(wav_filename, codec='pcm_s16le')
                        audio.close()
                        os.unlink(temp_file)
                    except Exception as e2:
                        # 如果都失败，直接返回mp3文件
                        os.rename(temp_file, wav_filename.replace('.wav', '.mp3'))
                        wav_filename = wav_filename.replace('.wav', '.mp3')
                
                if os.path.exists(wav_filename) and os.path.getsize(wav_filename) > 1024:
                    return True, f"gTTS生成成功，文件大小: {os.path.getsize(wav_filename)} 字节"
                else:
                    return False, "gTTS生成的文件无效"
                    
            except ImportError:
                return False, "gTTS模块未安装"
            except Exception as e:
                return False, f"gTTS生成失败: {str(e)}"
        
        def generate_audio_with_timeout(text, filename, rate, timeout=15):
            """带超时的音频生成函数 - 使用进程级保护"""
            try:
                # 清理文本
                cleaned_text = re.sub(r'[，。、；："\'\n\r]+', '，', text)
                cleaned_text = re.sub(r'[•●■◆]', '', cleaned_text)
                cleaned_text = re.sub(r'[，]+', '，', cleaned_text)
                cleaned_text = re.sub(r'^[，]+|[，]+$', '', cleaned_text)
                if not cleaned_text.strip():
                    cleaned_text = "空白幻灯片"
                if len(cleaned_text) > 100:
                    cleaned_text = cleaned_text[:97] + "..."
                
                # 使用更简单的方式生成临时脚本，避免JSON转义问题
                temp_script = tempfile.mktemp(suffix='.py')
                
                # 直接构建脚本内容
                with open(temp_script, 'w', encoding='utf-8') as f:
                    f.write('''
import pyttsx3
import os
import sys
import json

# 获取参数
text = ''' + json.dumps(cleaned_text) + '''
filename = ''' + json.dumps(filename) + '''
rate = ''' + str(rate or 150) + '''

result = {"success": False, "message": "未知错误"}

# 清理文本
cleaned_text = text
cleaned_text = cleaned_text.replace('\\n', '')

# 初始化引擎
try:
    engine = pyttsx3.init()
    engine.setProperty('rate', rate)
    
    # 设置语音
    voices = engine.getProperty('voices')
    for voice in voices:
        if any(keyword in voice.id.lower() for keyword in ['chinese', 'zh', 'mandarin', '中文', '普通话']):
            engine.setProperty('voice', voice.id)
            break
    
    # 生成音频
    engine.save_to_file(cleaned_text, filename)
    engine.runAndWait()
    
    # 验证文件
    if os.path.exists(filename) and os.path.getsize(filename) > 1024:
        result["success"] = True
        result["message"] = "生成成功"
    else:
        result["success"] = False
        result["message"] = "文件无效"
except Exception as e:
    result["success"] = False
    result["message"] = str(e)

# 输出结果
print(json.dumps(result, ensure_ascii=False))
''')
                
                try:
                    # 运行临时脚本
                    process = subprocess.Popen(
                        [sys.executable, temp_script],
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        encoding='utf-8',
                        errors='replace'
                    )
                    
                    # 等待超时
                    stdout, stderr = process.communicate(timeout=timeout)
                    
                    # 检查stdout是否为空
                    if not stdout:
                        # 尝试使用gTTS作为备选
                        self.log(f"    pyttsx3失败，尝试使用gTTS备选方案...")
                        return generate_audio_with_gtts(text, filename)
                            
                    # 确保输出是有效的JSON
                    try:
                        result = json.loads(stdout)
                        if result["success"]:
                            return result["success"], result["message"]
                        else:
                            # pyttsx3失败，尝试使用gTTS
                            self.log(f"    pyttsx3失败: {result['message']}，尝试使用gTTS备选方案...")
                            return generate_audio_with_gtts(text, filename)
                    except json.JSONDecodeError:
                        # JSON解析失败，尝试使用gTTS
                        self.log(f"    pyttsx3输出解析失败，尝试使用gTTS备选方案...")
                        return generate_audio_with_gtts(text, filename)
                        
                except subprocess.TimeoutExpired:
                    process.kill()
                    # 超时，尝试使用gTTS
                    self.log(f"    pyttsx3超时，尝试使用gTTS备选方案...")
                    return generate_audio_with_gtts(text, filename)
                except Exception as e:
                    # 执行错误，尝试使用gTTS
                    self.log(f"    pyttsx3执行错误: {str(e)}，尝试使用gTTS备选方案...")
                    return generate_audio_with_gtts(text, filename)
                finally:
                    # 清理临时脚本
                    if os.path.exists(temp_script):
                        try:
                            os.unlink(temp_script)
                        except:
                            pass
            except Exception as e:
                # 内部错误，尝试使用gTTS
                self.log(f"    音频生成内部错误: {str(e)}，尝试使用gTTS备选方案...")
                return generate_audio_with_gtts(text, filename)

        # 使用带超时的音频生成
        for i, text in enumerate(scripts):
            filename = os.path.join(temp_folder, f"audio_{i + 1:02d}.wav")
            self.log(f"  - 处理第 {i + 1} 页音频...")
            self.log(f"    原始文本: {text[:50]}..." if len(text) > 50 else f"    原始文本: {text}")
            self.log(f"    文本长度: {len(text)} 字符")

            try:
                # 检查gTTS是否可用
                self.log(f"    [调试] 检查gTTS模块是否可用...")
                try:
                    from gtts import gTTS
                    self.log(f"    [调试] ✓ gTTS模块已安装")
                    gtts_available = True
                except ImportError:
                    self.log(f"    [调试] ✗ gTTS模块未安装")
                    gtts_available = False
                    # 尝试安装gTTS
                    self.log(f"    [调试] 尝试安装gTTS...")
                    try:
                        import subprocess
                        install_result = subprocess.run(
                            [sys.executable, '-m', 'pip', 'install', 'gtts'],
                            capture_output=True,
                            text=True,
                            timeout=30
                        )
                        if install_result.returncode == 0:
                            self.log(f"    [调试] ✓ 成功安装gTTS")
                            gtts_available = True
                        else:
                            self.log(f"    [调试] ✗ 安装gTTS失败: {install_result.stderr[:100]}...")
                    except Exception as e:
                        self.log(f"    [调试] ✗ 安装gTTS异常: {e}")
                
                # 优先使用pyttsx3离线生成音频
                self.log(f"    开始生成语音...")
                success, message = generate_audio_with_timeout(text, filename, rate)
                
                # 如果pyttsx3失败，且gTTS可用，则尝试gTTS
                if not success and gtts_available:
                    self.log(f"    pyttsx3失败，尝试使用gTTS备选方案...")
                    success, message = generate_audio_with_gtts(text, filename)
                    if not success:
                        self.log(f"    gTTS也失败，尝试使用降级方案...")
                
                if success:
                    file_size = os.path.getsize(filename)
                    self.log(f"    ✓ 音频生成成功: {message}")
                    self.log(f"    ✓ 文件路径: {filename}")
                    self.log(f"    ✓ 文件大小: {file_size / 1024:.2f} KB")
                    audio_paths.append(filename)
                    self.log(f"    ✓ 音频已添加到列表")
                else:
                    self.log(f"    [警告] 首次生成失败: {message}")
                    
                    # 尝试重新生成一次
                    self.log(f"    尝试重新生成...")
                    success, message = generate_audio_with_timeout(text, filename, rate)
                    
                    if success:
                        file_size = os.path.getsize(filename)
                        self.log(f"    ✓ 重新生成成功: {message}")
                        self.log(f"    ✓ 文件大小: {file_size / 1024:.2f} KB")
                        audio_paths.append(filename)
                    else:
                        self.log(f"    [错误] 重新生成失败: {message}")
                        self.log(f"    尝试使用降级方案生成音频...")
                        # 降级方案
                        try:
                            from moviepy.audio.AudioClip import AudioArrayClip
                            estimated_duration = max(2.0, len(text) * 0.33)
                            self.log(f"    生成 {estimated_duration:.1f} 秒的音频...")
                            samples = int(estimated_duration * 44100)
                            t = np.linspace(0, estimated_duration, samples)
                            audio_array = np.sin(2 * np.pi * 440 * t)
                            audio_array = np.column_stack((audio_array, audio_array)).astype(np.float32)
                            audio_clip = AudioArrayClip(audio_array, fps=44100)
                            audio_clip.write_audiofile(filename, fps=44100)
                            
                            if os.path.exists(filename) and os.path.getsize(filename) > 1024:
                                file_size = os.path.getsize(filename)
                                self.log(f"    ✓ 降级方案生成成功")
                                self.log(f"    ✓ 文件大小: {file_size / 1024:.2f} KB")
                                audio_paths.append(filename)
                            else:
                                self.log(f"    [错误] 降级方案生成的文件无效")
                        except Exception as e3:
                            self.log(f"    [错误] 降级方案执行失败: {e3}")
            except Exception as e:
                self.log(f"    [错误] 音频生成过程发生异常: {e}")
                import traceback
                self.log(f"    异常详情: {traceback.format_exc()[:200]}...")
                # 降级方案
                self.log(f"    [降级方案] 使用简单音频生成替代")
                try:
                    from moviepy.audio.AudioClip import AudioArrayClip
                    estimated_duration = max(2.0, len(text) * 0.33)
                    samples = int(estimated_duration * 44100)
                    t = np.linspace(0, estimated_duration, samples)
                    audio_array = np.sin(2 * np.pi * 440 * t)
                    audio_array = np.column_stack((audio_array, audio_array)).astype(np.float32)
                    audio_clip = AudioArrayClip(audio_array, fps=44100)
                    audio_clip.write_audiofile(filename, fps=44100)
                    
                    if os.path.exists(filename) and os.path.getsize(filename) > 1024:
                        self.log(f"    降级方案生成成功")
                        audio_paths.append(filename)
                except Exception as e2:
                    self.log(f"    [错误] 降级方案也失败: {e2}")
                    continue
        
        self.log(f"  音频生成完成，共生成 {len(audio_paths)} 个音频文件")
        return audio_paths

    def make_video(self, images, audios, output_path, use_subtitle=False, scripts=None, timeout=None):
        # 导入必要的模块
        import os
        import sys
        # 添加直接文件写入，绕过self.log可能的问题
        def direct_log(message):
            # 直接写入绝对路径的日志文件
            log_path = os.path.abspath('ppt2video_direct.log')
            with open(log_path, 'a', encoding='utf-8') as f:
                f.write(f"{message}\n")
            # 同时调用原始log方法
            self.log(message)
        
        direct_log("=== 视频生成开始 (直接日志) ===")
        direct_log(f"当前时间: {__import__('datetime').datetime.now()}")
        direct_log(f"图片数量: {len(images)}")
        direct_log(f"音频数量: {len(audios)}")
        direct_log(f"输出路径: {output_path}")
        direct_log(f"当前工作目录: {os.getcwd()}")
        direct_log(f"Python版本: {sys.version}")
        
        self.log("正在进行视频渲染 (这可能需要几分钟)...")
        clips = []
        
        # 记录关键调试信息
        self.log(f"=== 视频生成核心参数 ===")
        self.log(f"  图片数量: {len(images)}")
        self.log(f"  音频数量: {len(audios)}")
        self.log(f"  输出路径: {output_path}")
        self.log(f"  当前工作目录: {os.getcwd()}")
        self.log(f"  Python版本: {sys.version}")
        
        # 检查FFmpeg配置
        self.log(f"=== FFmpeg配置检查 ===")
        try:
            import imageio_ffmpeg
            self.log(f"  imageio_ffmpeg版本: {imageio_ffmpeg.__version__}")
            ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
            self.log(f"  FFmpeg路径: {ffmpeg_path}")
            self.log(f"  FFmpeg存在: {os.path.exists(ffmpeg_path)}")
            if os.path.exists(ffmpeg_path):
                self.log(f"  FFmpeg文件大小: {os.path.getsize(ffmpeg_path)} 字节")
        except Exception as e:
            self.log(f"  FFmpeg检查失败: {e}")

        # 确保输出目录存在
        output_dir = os.path.dirname(output_path)
        self.log(f"=== 输出目录处理 ===")
        self.log(f"  输出目录: {output_dir}")
        self.log(f"  输出目录存在: {os.path.exists(output_dir)}")
        
        if output_dir and not os.path.exists(output_dir):
            self.log(f"  创建输出目录: {output_dir}")
            try:
                os.makedirs(output_dir)
                self.log(f"  输出目录创建成功")
            except Exception as e:
                self.log(f"  创建输出目录失败: {e}")

        # 确保输出路径是绝对路径
        output_path = os.path.abspath(output_path)
        self.log(f"  绝对输出路径: {output_path}")

        # 验证输入的图片和音频列表
        self.log(f"=== 输入文件验证 ===")
        self.log(f"  图片列表: {[os.path.basename(img) for img in images]}")
        self.log(f"  音频列表: {[os.path.basename(aud) for aud in audios]}")
        
        # 检查输入文件是否存在
        valid_pairs = []
        for i, (img, aud) in enumerate(zip(images, audios)):
            img_exists = os.path.exists(img)
            aud_exists = os.path.exists(aud)
            img_size = os.path.getsize(img) if img_exists else 0
            aud_size = os.path.getsize(aud) if aud_exists else 0
            self.log(f"  第{i+1}对文件:")
            self.log(f"    图片: {os.path.basename(img)} - 存在: {img_exists}, 大小: {img_size} 字节")
            self.log(f"    音频: {os.path.basename(aud)} - 存在: {aud_exists}, 大小: {aud_size} 字节")
            
            if img_exists and aud_exists:
                if img_size < 1024 or aud_size < 1024:
                    self.log(f"    [警告] 文件太小")
                valid_pairs.append((img, aud))
        
        self.log(f"  有效图片-音频对数量: {len(valid_pairs)}")
        
        # 处理每一对图片和音频
        for i, (img, aud) in enumerate(valid_pairs):
            try:
                self.log(f"  处理第 {i + 1} 页:")
                self.log(f"    图片路径: {img}")
                self.log(f"    音频路径: {aud}")
                
                # 检查文件大小
                img_size = os.path.getsize(img)
                aud_size = os.path.getsize(aud)
                self.log(f"    图片大小: {img_size} 字节")
                self.log(f"    音频大小: {aud_size} 字节")
                
                if img_size < 1024:  # 小于1KB的图片可能有问题
                    self.log(f"    [警告] 图片文件太小，跳过此页")
                    continue
                
                if aud_size < 1024:  # 小于1KB的音频可能有问题
                    self.log(f"    [警告] 音频文件太小，跳过此页")
                    continue

                # 加载音频
                self.log(f"    加载音频...")
                try:
                    audio_clip = AudioFileClip(aud)
                    duration = audio_clip.duration  # 使用音频的实际时长作为幻灯片时长，确保音频能完整播放
                    self.log(f"    音频时长: {duration} 秒，幻灯片将显示相同时长")
                    self.log(f"    音频采样率: {audio_clip.fps}")
                    self.log(f"    音频通道数: {audio_clip.nchannels}")
                except Exception as e:
                    self.log(f"    [错误] 加载音频失败: {e}")
                    # 不再创建测试音频，直接跳过此页
                    continue

                # 加载图片
                self.log(f"    加载图片...")
                try:
                    img_clip = ImageClip(img).with_duration(duration)
                except Exception as e:
                    self.log(f"    [错误] 加载图片失败: {e}")
                    continue

                # 转场效果
                if i > 0:
                    self.log(f"    添加转场效果...")
                    try:
                        img_clip = img_clip.crossfadein(0.8)
                    except Exception as e:
                        self.log(f"    [警告] 添加转场效果失败(跳过): {e}")

                final_clip = img_clip

                # 字幕 (简易版)
                if use_subtitle and scripts:
                    # 注意：如果没装 ImageMagick，这里会报错，所以默认建议关闭
                    try:
                        self.log(f"    生成字幕...")
                        txt = scripts[i][:30]  # 防止太长
                        txt_clip = TextClip(txt, fontsize=40, color='white', bg_color='rgba(0,0,0,0.6)',
                                            font='Microsoft-YaHei-Bold', size=(img_clip.w * 0.9, None),
                                            method='caption')
                        txt_clip = txt_clip.with_pos(('center', 'bottom')).with_duration(duration)
                        final_clip = CompositeVideoClip([img_clip, txt_clip])
                    except Exception as e:
                        self.log(f"    [警告] 字幕生成失败(跳过): {e}")

                self.log(f"    组合音视频片段...")
                try:
                    final_clip = final_clip.with_audio(audio_clip)
                    # 验证音频是否成功附加
                    if hasattr(final_clip, 'audio') and final_clip.audio is not None:
                        self.log(f"    ✓ 音频成功附加到视频")
                        self.log(f"    附加后的音频时长: {final_clip.audio.duration}")
                    else:
                        self.log(f"    ✗ 音频未能附加到视频")
                        continue
                    clips.append(final_clip)
                    self.log(f"    第 {i + 1} 页处理完成")
                except Exception as e:
                    self.log(f"    [错误] 组合音视频片段失败: {e}")
                    import traceback
                    traceback.print_exc()
                    continue
            except Exception as e:
                self.log(f"    [错误] 处理第 {i + 1} 页时出错: {e}")
                import traceback
                traceback.print_exc()

        self.log(f"  总共有 {len(clips)} 个有效片段")
        direct_log(f"  总共有 {len(clips)} 个有效片段")
        
        if not clips:
            error_msg = "生成失败：没有有效的视频片段。请检查PPT文件是否有内容。"
            self.log(error_msg)
            direct_log(error_msg)
            # 不再创建测试视频，直接返回错误
            raise Exception(error_msg)
        
        self.log(f"  合并视频片段...")
        direct_log(f"  合并视频片段...")
        
        # 验证每个片段是否有音频
        for i, clip in enumerate(clips):
            has_clip_audio = hasattr(clip, 'audio') and clip.audio is not None
            direct_log(f"  片段{i+1}是否有音频: {has_clip_audio}")
            if has_clip_audio:
                direct_log(f"    片段音频时长: {clip.audio.duration}")
        
        try:
            final_video = concatenate_videoclips(clips, method="compose")
            self.log(f"  合并完成，开始编码视频...")
            direct_log(f"  合并完成，开始编码视频...")
            direct_log(f"  合并后的视频时长: {final_video.duration if hasattr(final_video, 'duration') else '未知'}")
            direct_log(f"  视频是否有音频: {hasattr(final_video, 'audio') and final_video.audio is not None}")
            
            # 如果合并后没有音频，尝试手动添加音频
            if not hasattr(final_video, 'audio') or final_video.audio is None:
                direct_log(f"  警告：合并后视频没有音频，尝试手动合并音频")
                from moviepy.audio.AudioClip import concatenate_audioclips
                
                # 收集所有片段的音频
                audio_clips = [clip.audio for clip in clips if hasattr(clip, 'audio') and clip.audio is not None]
                if audio_clips:
                    direct_log(f"  收集到 {len(audio_clips)} 个音频片段")
                    # 合并音频
                    final_audio = concatenate_audioclips(audio_clips)
                    # 附加到视频
                    final_video = final_video.with_audio(final_audio)
                    direct_log(f"  手动合并音频成功")
                    direct_log(f"  手动合并后视频是否有音频: {hasattr(final_video, 'audio') and final_video.audio is not None}")
            
            # 验证输出目录存在
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
                self.log(f"  再次确认创建输出目录: {output_dir}")
            
            # 确保输出路径是绝对路径
            output_path = os.path.abspath(output_path)
            self.log(f"  绝对输出路径: {output_path}")
            
            # 添加更多编码参数，确保视频文件完整可播放
            try:
                # 尝试使用更简单的参数设置，确保视频能正常生成
                self.log(f"=== 开始视频编码 ===")
                self.log(f"  输出路径: {output_path}")
                self.log(f"  编码参数:")
                self.log(f"    FPS: 24")
                self.log(f"    视频编码: libx264")
                self.log(f"    音频编码: aac")
                self.log(f"    预设: ultrafast")
                self.log(f"    线程数: 1")
                self.log(f"    视频比特率: 1000k")
                self.log(f"    音频比特率: 64k")
                self.log(f"    像素格式: yuv420p")
                
                # 添加直接日志
                direct_log("=== 开始视频编码 (直接日志) ===")
                direct_log(f"  输出路径: {output_path}")
                direct_log(f"  开始编码时间: {__import__('datetime').datetime.now()}")
                
                # 移除了测试视频生成代码，直接进行正式视频编码
                
                # 使用write_videofile编码主视频
                direct_log("\n--- 编码主视频 ---")
                direct_log(f"  视频时长: {final_video.duration if hasattr(final_video, 'duration') else '未知'}")
                direct_log(f"  视频是否有音频: {hasattr(final_video, 'audio') and final_video.audio is not None}")
                direct_log(f"  输出路径: {output_path}")
                
                # 使用更简单的参数，禁用所有复杂功能
                try:
                    # 验证输出目录存在
                    output_dir = os.path.dirname(output_path)
                    if output_dir and not os.path.exists(output_dir):
                        os.makedirs(output_dir)
                        direct_log(f"  创建输出目录: {output_dir}")
                    
                    # 验证输出路径的可写性
                    try:
                        with open(output_path, 'w') as test_file:
                            test_file.write('test')
                        os.remove(output_path)
                        direct_log(f"  输出路径可写性测试通过")
                    except Exception as write_test_e:
                        direct_log(f"  [错误] 输出路径可写性测试失败: {write_test_e}")
                        raise Exception(f"输出路径不可写: {output_path}")
                    
                    # 创建自定义进度反馈函数
                    def progress_callback(progress):
                        percent = progress * 100
                        direct_log(f"  编码进度: {percent:.1f}% - {progress:.3f}")
                    
                    # 视频编码的结果和异常
                    encode_result = None
                    encode_exception = None
                    
                    # 定义编码函数
                    def encode_video():
                        nonlocal encode_result, encode_exception
                        try:
                            direct_log("  启动编码线程...")
                            # 再次验证视频是否有音频
                            has_audio = hasattr(final_video, 'audio') and final_video.audio is not None
                            direct_log(f"  编码前视频是否有音频: {has_audio}")
                            if has_audio:
                                direct_log(f"  音频时长: {final_video.audio.duration}")
                                direct_log(f"  音频采样率: {final_video.audio.fps}")
                                direct_log(f"  音频通道数: {final_video.audio.nchannels}")
                            
                            # 添加更详细的音频验证
                            has_audio = hasattr(final_video, 'audio') and final_video.audio is not None
                            direct_log(f"  最终视频是否有音频: {has_audio}")
                            if has_audio:
                                direct_log(f"  音频时长: {final_video.audio.duration}")
                                direct_log(f"  音频采样率: {final_video.audio.fps}")
                                direct_log(f"  音频通道数: {final_video.audio.nchannels}")
                            
                            # 修改ffmpeg参数，使用更快的编码设置减少卡住的可能性
                            final_video.write_videofile(
                                output_path,
                                fps=8,  # 降低帧率以提高编码速度
                                codec="libx264",
                                audio_codec="aac",  # 始终使用aac音频编码
                                preset="ultrafast",  # 使用最快的预设以减少转换时间
                                threads=4,  # 增加线程数利用多核CPU
                                bitrate="500k",  # 合理的视频比特率平衡质量和速度
                                ffmpeg_params=[
                                    '-pix_fmt', 'yuv420p',  # 使用兼容性更好的像素格式
                                    '-b:a', '96k',  # 适当降低音频比特率
                                    '-y',  # 覆盖现有文件
                                    '-loglevel', 'error',  # 减少日志输出提高性能
                                    '-ac', '1',  # 使用单声道减少音频处理量
                                    '-ar', '44100',  # 保持44.1kHz采样率
                                    '-movflags', '+faststart',  # 优化视频以便快速开始播放
                                    '-profile:v', 'baseline',  # 使用基线配置文件，提高兼容性
                                    '-level', '3.0'  # 编码级别，确保广泛兼容
                                ]
                            )
                            encode_result = True
                            direct_log("  编码线程完成")
                            
                            # 验证生成的视频文件是否包含音频
                            if os.path.exists(output_path):
                                try:
                                    import subprocess
                                    direct_log("  使用ffprobe检查视频音频流...")
                                    # 检查视频中的音频流
                                    result = subprocess.run(
                                        ['ffprobe', '-v', 'error', '-select_streams', 'a', '-show_entries', 'stream=codec_name', '-of', 'csv=p=0', output_path],
                                        capture_output=True,
                                        text=True
                                    )
                                    if result.stdout.strip():
                                        direct_log(f"  ✓ 生成的视频包含音频流: {result.stdout.strip()}")
                                    else:
                                        direct_log("  ✗ 生成的视频不包含音频流")
                                        
                                    # 检查视频总流数
                                    result = subprocess.run(
                                        ['ffprobe', '-v', 'error', '-show_entries', 'format=nb_streams', '-of', 'csv=p=0', output_path],
                                        capture_output=True,
                                        text=True
                                    )
                                    if result.stdout.strip():
                                        direct_log(f"  视频总流数: {result.stdout.strip()}")
                                except Exception as ffprobe_e:
                                    direct_log(f"  ffprobe检查失败: {ffprobe_e}")
                        except Exception as e:
                            encode_exception = e
                            encode_result = False
                            direct_log(f"  编码线程异常: {e}")
                            import traceback
                            traceback.print_exc()
                    
                    # 计算超时时间（视频时长的10倍，至少30秒，最多30分钟）
                    video_duration = final_video.duration if hasattr(final_video, 'duration') else 10
                    # 如果传递了timeout参数，就使用该参数，否则按照原有逻辑计算
                    if timeout is None:
                        timeout_seconds = max(30, min(1800, video_duration * 10))
                    else:
                        timeout_seconds = timeout
                    direct_log(f"  设置编码超时时间: {timeout_seconds}秒")
                    
                    import threading
                    # 启动编码线程
                    encode_thread = threading.Thread(target=encode_video)
                    encode_thread.start()
                    
                    # 等待编码完成或超时
                    encode_thread.join(timeout_seconds)
                    
                    # 检查是否超时
                    if encode_thread.is_alive():
                        # 编码超时
                        direct_log(f"[错误] 视频编码超时 ({timeout_seconds}秒)")
                        # 尝试终止线程并清理资源
                        try:
                            # 记录超时错误
                            with open('ppt2video_timeout.log', 'a', encoding='utf-8') as f:
                                f.write(f"=== 编码超时 ===")
                                f.write(f"时间: {__import__('datetime').datetime.now()}")
                                f.write(f"视频时长: {video_duration}秒")
                                f.write(f"超时时间: {timeout_seconds}秒")
                                f.write(f"\n")
                        except Exception as e:
                            direct_log(f"[错误] 记录超时信息失败: {e}")
                            
                        raise Exception(f"视频编码超时 ({timeout_seconds}秒)")
                    
                    # 检查编码结果
                    if encode_result:
                        direct_log("  编码线程执行成功")
                    else:
                        direct_log(f"  编码线程执行失败: {encode_exception}")
                        raise encode_exception
                    
                except Exception as e:
                    direct_log(f"  [错误] 视频编码失败: {e}")
                    import traceback
                    traceback.print_exc()
                    raise
            except Exception as e:
                direct_log(f"  [错误] 视频编码过程中的异常: {e}")
                import traceback
                traceback.print_exc()
                raise
        except Exception as e:
            direct_log(f"  [错误] 视频生成失败: {e}")
            import traceback
            traceback.print_exc()
            raise

        self.log(f"视频生成完成，输出路径: {output_path}")
        direct_log(f"视频生成完成，输出路径: {output_path}")
        return output_path

# ===========================
# GUI界面类
# ===========================
class PPTtoVideoConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("PPT转视频工具")
        self.root.geometry("800x600")
        
        # 设置中文字体
        self.set_chinese_font()
        
        # 创建转换逻辑实例
        self.converter = ConverterLogic(self.log)
        
        # 界面元素
        self.create_widgets()
        
        # 进度信息
        self.is_converting = False

    def set_chinese_font(self):
        """设置中文字体，解决Windows下中文显示问题"""
        try:
            # Windows系统
            if sys.platform == 'win32':
                # 使用系统默认中文字体
                from tkinter import font
                default_font = font.nametofont("TkDefaultFont")
                default_font.configure(family="SimHei", size=10)
                self.root.option_add("*Font", default_font)
        except Exception as e:
            print(f"设置中文字体失败: {e}")

    def create_widgets(self):
        # 创建主框架
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 标题
        title_frame = ttk.Frame(main_frame)
        title_frame.pack(fill=tk.X, pady=10)
        ttk.Label(title_frame, text="PPT转视频工具", font=('SimHei', 16)).pack(side=tk.LEFT)
        
        # 输入区
        input_frame = ttk.LabelFrame(main_frame, text="输入设置", padding="10")
        input_frame.pack(fill=tk.X, pady=10)
        
        # PPT文件选择
        ppt_frame = ttk.Frame(input_frame)
        ppt_frame.pack(fill=tk.X, pady=5)
        ttk.Label(ppt_frame, text="PPT文件:", width=10).pack(side=tk.LEFT)
        self.ppt_path_var = tk.StringVar()
        ttk.Entry(ppt_frame, textvariable=self.ppt_path_var, width=50).pack(side=tk.LEFT, padx=5)
        ttk.Button(ppt_frame, text="浏览", command=self.browse_ppt).pack(side=tk.LEFT)
        
        # 输出路径
        output_frame = ttk.Frame(input_frame)
        output_frame.pack(fill=tk.X, pady=5)
        ttk.Label(output_frame, text="输出路径:", width=10).pack(side=tk.LEFT)
        self.output_path_var = tk.StringVar()
        ttk.Entry(output_frame, textvariable=self.output_path_var, width=50).pack(side=tk.LEFT, padx=5)
        ttk.Button(output_frame, text="浏览", command=self.browse_output).pack(side=tk.LEFT)
        
        # 选项区
        options_frame = ttk.LabelFrame(main_frame, text="高级选项", padding="10")
        options_frame.pack(fill=tk.X, pady=10)
        
        # 语音选择
        voice_frame = ttk.Frame(options_frame)
        voice_frame.pack(fill=tk.X, pady=5)
        ttk.Label(voice_frame, text="语音选择:", width=10).pack(side=tk.LEFT)
        self.voice_var = tk.StringVar()
        ttk.Combobox(voice_frame, textvariable=self.voice_var, values=["中文", "英文"]).pack(side=tk.LEFT, padx=5)
        self.voice_var.set("中文")
        
        # 语速设置
        rate_frame = ttk.Frame(options_frame)
        rate_frame.pack(fill=tk.X, pady=5)
        ttk.Label(rate_frame, text="语速设置:", width=10).pack(side=tk.LEFT)
        self.rate_var = tk.IntVar(value=150)
        ttk.Scale(rate_frame, from_=50, to=250, variable=self.rate_var, orient=tk.HORIZONTAL).pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        ttk.Label(rate_frame, textvariable=self.rate_var, width=5).pack(side=tk.LEFT)
        
        # 字幕选项
        subtitle_frame = ttk.Frame(options_frame)
        subtitle_frame.pack(fill=tk.X, pady=5)
        self.subtitle_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(subtitle_frame, text="生成字幕 (需安装ImageMagick)", variable=self.subtitle_var).pack(side=tk.LEFT)
        
        # 转换按钮
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill=tk.X, pady=10)
        ttk.Button(button_frame, text="开始转换", command=self.start_conversion, state=tk.NORMAL).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="停止转换", command=self.stop_conversion, state=tk.DISABLED).pack(side=tk.LEFT, padx=5)
        
        # 日志区
        log_frame = ttk.LabelFrame(main_frame, text="转换日志", padding="10")
        log_frame.pack(fill=tk.BOTH, expand=True, pady=10)
        self.log_text = scrolledtext.ScrolledText(log_frame, width=80, height=20)
        self.log_text.pack(fill=tk.BOTH, expand=True)

    def browse_ppt(self):
        file_path = filedialog.askopenfilename(
            title="选择PPT文件",
            filetypes=[("PPT文件", "*.ppt *.pptx"), ("所有文件", "*.*")]
        )
        if file_path:
            self.ppt_path_var.set(file_path)
            # 自动设置输出路径
            base_name = os.path.splitext(file_path)[0]
            self.output_path_var.set(base_name + ".mp4")

    def browse_output(self):
        file_path = filedialog.asksaveasfilename(
            title="选择输出文件",
            filetypes=[("MP4视频", "*.mp4"), ("所有文件", "*.*")],
            defaultextension=".mp4"
        )
        if file_path:
            self.output_path_var.set(file_path)

    def log(self, message):
        """日志输出到界面"""
        current_time = __import__('datetime').datetime.now().strftime("%H:%M:%S")
        log_message = f"[{current_time}] {message}\n"
        self.log_text.insert(tk.END, log_message)
        self.log_text.see(tk.END)
        self.log_text.update_idletasks()

    def start_conversion(self):
        """开始转换线程"""
        if self.is_converting:
            return
        
        # 获取输入
        ppt_path = self.ppt_path_var.get().strip()
        output_path = self.output_path_var.get().strip()
        voice_id = "chinese" if self.voice_var.get() == "中文" else "english"
        rate = self.rate_var.get()
        use_subtitle = self.subtitle_var.get()
        
        # 验证输入
        if not ppt_path or not os.path.exists(ppt_path):
            messagebox.showerror("错误", "请选择有效的PPT文件")
            return
        
        if not output_path:
            messagebox.showerror("错误", "请选择输出路径")
            return
        
        # 清空日志
        self.log_text.delete(1.0, tk.END)
        self.log("=== 开始转换 ===")
        self.log(f"PPT文件: {ppt_path}")
        self.log(f"输出路径: {output_path}")
        self.log(f"语音类型: {self.voice_var.get()}")
        self.log(f"语速: {rate}")
        self.log(f"生成字幕: {'是' if use_subtitle else '否'}")
        
        # 禁用按钮
        self.is_converting = True
        self.root.nametowidget(".!frame2.!frame3.!button").configure(state=tk.DISABLED)
        self.root.nametowidget(".!frame2.!frame3.!button2").configure(state=tk.NORMAL)
        
        # 开始转换线程
        self.conversion_thread = threading.Thread(
            target=self.convert_ppt_to_video,
            args=(ppt_path, output_path, voice_id, rate, use_subtitle)
        )
        self.conversion_thread.daemon = True
        self.conversion_thread.start()

    def stop_conversion(self):
        """停止转换"""
        self.is_converting = False
        self.log("停止转换中...")

    def convert_ppt_to_video(self, ppt_path, output_path, voice_id, rate, use_subtitle):
        """转换PPT到视频的核心逻辑"""
        try:
            pythoncom.CoInitialize()  # <--- 必须加这个，解决COM对象初始化问题
            
            # 临时目录
            temp_folder = os.path.join(os.path.dirname(ppt_path), "temp_ppt_converter")
            if os.path.exists(temp_folder):
                shutil.rmtree(temp_folder)
            os.makedirs(temp_folder)
            
            self.log(f"临时目录: {temp_folder}")
            
            try:
                # 1. 导出PPT为图片
                self.log("\n1. 导出PPT为图片")
                images = self.converter.export_images(ppt_path, temp_folder)
                if not images:
                    raise Exception("导出图片失败")
                self.log(f"成功导出 {len(images)} 张图片")
                
                # 2. 提取文字脚本
                self.log("\n2. 提取文字脚本")
                scripts = self.converter.extract_text(ppt_path)
                if not scripts:
                    self.log("[警告] 未提取到文字内容，使用默认文本")
                    scripts = [f"第 {i+1} 页" for i in range(len(images))]
                self.log(f"成功提取 {len(scripts)} 段文字")
                
                # 3. 生成音频
                self.log("\n3. 生成音频")
                audios = self.converter.generate_audio(scripts, temp_folder, voice_id, rate)
                if not audios:
                    raise Exception("生成音频失败")
                self.log(f"成功生成 {len(audios)} 个音频文件")
                
                # 4. 生成视频
                self.log("\n4. 生成视频")
                video_path = self.converter.make_video(images, audios, output_path, use_subtitle, scripts)
                self.log(f"视频生成成功: {video_path}")
                
                # 转换完成
                self.log("\n=== 转换完成 ===")
                messagebox.showinfo("转换完成", f"视频已成功生成:\n{video_path}")
                
            except Exception as e:
                self.log(f"\n[错误] 转换失败: {e}")
                import traceback
                traceback.print_exc()
                messagebox.showerror("转换失败", f"转换过程中发生错误:\n{e}")
            finally:
                # 清理临时目录
                if os.path.exists(temp_folder):
                    try:
                        shutil.rmtree(temp_folder)
                        self.log(f"临时目录已清理: {temp_folder}")
                    except Exception as e:
                        self.log(f"清理临时目录失败: {e}")
                
        except Exception as e:
            self.log(f"\n[致命错误] {e}")
            import traceback
            traceback.print_exc()
        finally:
            # 启用按钮
            self.is_converting = False
            self.root.nametowidget(".!frame2.!frame3.!button").configure(state=tk.NORMAL)
            self.root.nametowidget(".!frame2.!frame3.!button2").configure(state=tk.DISABLED)
            
            # 释放COM资源
            try:
                pythoncom.CoUninitialize()
            except:
                pass

# 主程序入口
if __name__ == "__main__":
    root = tk.Tk()
    app = PPTtoVideoConverter(root)
    root.mainloop()


# ===========================
# 对外接口函数
# ===========================
def convert_presentation_to_video(ppt_path, output_path, rate=150, timeout=None):
    """
    转换PPT到视频的对外接口函数
    
    参数:
        ppt_path: PPT文件路径
        output_path: 输出视频路径
        rate: 语速 (默认150)
        
    返回:
        bool: 转换是否成功
    """
    import pythoncom
    import shutil
    import os
    
    # 初始化COM对象
    pythoncom.CoInitialize()
    
    # 简单的日志函数
    def simple_log(message):
        print(f"[PPT2VIDEO] {message}")
    
    try:
        # 创建转换器实例
        converter = ConverterLogic(simple_log)
        
        # 临时目录
        temp_folder = os.path.join(os.path.dirname(ppt_path), "temp_ppt_converter")
        if os.path.exists(temp_folder):
            shutil.rmtree(temp_folder)
        os.makedirs(temp_folder)
        
        simple_log(f"临时目录: {temp_folder}")
        
        # 1. 导出PPT为图片
        simple_log("1. 导出PPT为图片")
        images = converter.export_images(ppt_path, temp_folder)
        if not images:
            simple_log("导出图片失败")
            return False
        simple_log(f"成功导出 {len(images)} 张图片")
        
        # 2. 提取文字脚本
        simple_log("2. 提取文字脚本")
        scripts = converter.extract_text(ppt_path)
        if not scripts:
            simple_log("[警告] 未提取到文字内容，使用默认文本")
            scripts = [f"第 {i+1} 页" for i in range(len(images))]
        simple_log(f"成功提取 {len(scripts)} 段文字")
        
        # 3. 生成音频
        simple_log("3. 生成音频")
        audios = converter.generate_audio(scripts, temp_folder, None, rate)
        if not audios:
            simple_log("生成音频失败")
            return False
        simple_log(f"成功生成 {len(audios)} 个音频文件")
        
        # 4. 生成视频
        simple_log("4. 生成视频")
        video_path = converter.make_video(images, audios, output_path, False, scripts, timeout=timeout)
        simple_log(f"视频生成成功: {video_path}")
        
        return True
        
    except Exception as e:
        simple_log(f"转换失败: {e}")
        import traceback
        traceback.print_exc()
        return False
    finally:
        # 清理临时目录
        if 'temp_folder' in locals() and os.path.exists(temp_folder):
            try:
                shutil.rmtree(temp_folder)
                simple_log(f"已清理临时目录: {temp_folder}")
            except:
                pass
        
        # 释放COM资源
        try:
            pythoncom.CoUninitialize()
        except:
            pass