import os
import sys
import logging
import tempfile
from datetime import datetime

# 设置日志
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('direct_test_detailed.log', encoding='utf-8'),
        logging.StreamHandler(sys.stdout)
    ]
)

logger = logging.getLogger()

# 获取当前脚本目录，确保能导入ppt2video
sys.path.append(os.path.abspath('.'))

# 导入核心函数
try:
    from ppt2video import convert_presentation_to_video
    logger.info("成功导入convert_presentation_to_video函数")
except ImportError as e:
    logger.error(f"导入失败: {e}")
    sys.exit(1)

# 创建一个测试PPTX文件（如果没有的话）
def create_test_ppt(output_path):
    """创建一个简单的测试PPTX文件"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        # 创建演示文稿
        prs = Presentation()

        # 创建幻灯片1：标题幻灯片
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]

        title.text = "测试演示文稿"
        subtitle.text = "用于PPT转视频测试"

        # 设置标题颜色和大小
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title.text_frame.paragraphs[0].font.size = Pt(32)

        # 创建幻灯片2：内容幻灯片
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title2 = slide2.shapes.title
        content = slide2.placeholders[1]

        title2.text = "幻灯片2"
        content.text = "这是第二张幻灯片的内容。\n用于测试PPT转视频功能。"

        # 创建幻灯片3：带图片的幻灯片
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        title3 = slide3.shapes.title
        title3.text = "幻灯片3"

        # 添加一些文本框
        left = top = width = height = Inches(1)
        txBox = slide3.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "第三张幻灯片"

        # 添加备注
        notes_slide = slide3.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = "这是第三张幻灯片的备注内容。"

        # 保存演示文稿
        prs.save(output_path)
        logger.info(f"成功创建测试PPT: {output_path}")
        return True
    except Exception as e:
        logger.error(f"创建测试PPT失败: {e}")
        return False

# 主测试函数
def main():
    logger.info("===== 开始PPT转视频详细测试 =====")
    logger.info(f"Python版本: {sys.version}")
    logger.info(f"当前工作目录: {os.getcwd()}")

    # 创建临时目录
    temp_dir = tempfile.mkdtemp(prefix="ppt2video_test_")
    logger.info(f"创建临时目录: {temp_dir}")

    try:
        # 创建测试PPT
        ppt_path = os.path.join(temp_dir, "test_presentation.pptx")
        if not os.path.exists(ppt_path):
            if not create_test_ppt(ppt_path):
                logger.error("无法创建测试PPT，终止测试")
                return False

        # 验证PPT存在
        if os.path.exists(ppt_path):
            logger.info(f"测试PPT存在，大小: {os.path.getsize(ppt_path)} 字节")
        else:
            logger.error(f"测试PPT不存在: {ppt_path}")
            return False

        # 输出路径
        output_path = os.path.join(temp_dir, "output_test.mp4")
        logger.info(f"输出视频路径: {output_path}")

        # 调用转换函数
        logger.info("===== 开始转换PPT到视频 =====")
        start_time = datetime.now()
        
        # 直接调用核心函数
        success = convert_presentation_to_video(ppt_path, output_path)
        
        end_time = datetime.now()
        elapsed = end_time - start_time
        logger.info(f"转换完成，耗时: {elapsed.total_seconds():.2f} 秒")
        logger.info(f"转换结果: {'成功' if success else '失败'}")

        # 检查输出文件
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"输出视频文件存在")
            logger.info(f"文件路径: {output_path}")
            logger.info(f"文件大小: {file_size} 字节")
            logger.info(f"文件大小: {file_size / 1024:.2f} KB")
            logger.info(f"文件大小: {file_size / (1024 * 1024):.2f} MB")
            
            if file_size < 1024:
                logger.error(f"❌ 错误: 视频文件太小，只有 {file_size} 字节")
            else:
                logger.info(f"✅ 视频文件大小正常")
        else:
            logger.error(f"❌ 错误: 输出视频文件不存在")

        return success

    except Exception as e:
        logger.error(f"测试过程中出现异常: {e}", exc_info=True)
        return False
    finally:
        logger.info(f"===== 测试结束 =====")

if __name__ == "__main__":
    main()