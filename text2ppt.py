import tkinter as tk
from tkinter import ttk, scrolledtext, filedialog, messagebox
import threading
import os
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ===========================
# 1. 核心配置与主题
# ===========================
THEMES = {
    "商务蓝 (Business Blue)": {
        "bg_color": RGBColor(245, 247, 250),
        "title_color": RGBColor(23, 43, 77),
        "text_color": RGBColor(66, 82, 110),
        "accent_color": RGBColor(0, 82, 204),
        "highlight_color": RGBColor(255, 86, 48),
        "font_name": "Microsoft YaHei"
    },
    "莫兰迪 (Morandi)": {
        "bg_color": RGBColor(245, 245, 247),
        "title_color": RGBColor(44, 62, 80),
        "text_color": RGBColor(90, 95, 105),
        "accent_color": RGBColor(118, 146, 166),
        "highlight_color": RGBColor(214, 137, 115),
        "font_name": "Microsoft YaHei"
    },
    "暗黑科技 (Dark Tech)": {
        "bg_color": RGBColor(30, 30, 35),
        "title_color": RGBColor(235, 235, 235),
        "text_color": RGBColor(170, 175, 180),
        "accent_color": RGBColor(94, 189, 178),
        "highlight_color": RGBColor(255, 179, 71),
        "font_name": "Microsoft YaHei"
    }
}


# ===========================
# 2. 逻辑模块：文本解析与图片下载
# ===========================
class SimpleTextAnalyzer:
    """
    简单的文本解析器。
    规则：
    1. 全文第一行 = 封面标题
    2. 空行分隔不同页面
    3. 每个块的第一行 = 页面标题，后续行 = 要点
    """

    def __init__(self, raw_text):
        self.raw_text = raw_text

    def analyze(self):
        # 统一换行符并去除前后空白
        text = self.raw_text.strip().replace('\r\n', '\n')
        if not text:
            return []

        blocks = text.split('\n\n')  # 按空行分割
        slides_data = []

        # 处理封面 (第一个块)
        title_block = blocks[0].strip().split('\n')
        main_title = title_block[0]
        sub_title = title_block[1] if len(title_block) > 1 else "自动生成演示文稿"

        slides_data.append({
            "type": "title",
            "title": main_title,
            "subtitle": sub_title
        })

        # 处理内容页 (后续块)
        for block in blocks[1:]:
            lines = [l.strip() for l in block.split('\n') if l.strip()]
            if not lines:
                continue

            page_title = lines[0]
            points = lines[1:]

            # 简单提取关键词用于配图 (取标题的前两个词，实际可用 Jieba 分词优化)
            # 这里简单用标题作为搜索词
            keyword = page_title

            slides_data.append({
                "type": "content",
                "title": page_title,
                "points": points,
                "image_keyword": keyword
            })

        # 添加结束页
        slides_data.append({
            "type": "title",
            "title": "谢谢观看",
            "subtitle": "Q & A"
        })

        return slides_data


def get_image_stream(keyword):
    """下载图片流"""
    try:
        # 简单的处理中文关键词转 URL 编码，避免乱码
        # 使用 pollinations.ai (无需 Key，支持中文 Prompt)
        print(f"正在下载配图: {keyword}")
        url = f"https://image.pollinations.ai/prompt/{keyword}?width=800&height=600&nologo=true"
        response = requests.get(url, timeout=5)
        if response.status_code == 200:
            return BytesIO(response.content)
    except Exception as e:
        print(f"配图下载失败: {e}")
    return None


# ===========================
# 3. 核心模块：PPT 生成器 (复用之前的高级版)
# ===========================
class PPTGenerator:
    # 主题映射表：将前端的简写键映射到后端的完整键名
    THEME_MAPPING = {
        "business_blue": "商务蓝 (Business Blue)",
        "morandi": "莫兰迪 (Morandi)",
        "dark_tech": "暗黑科技 (Dark Tech)"
    }
    
    def __init__(self, filename, theme_key):
        self.prs = Presentation()
        self.filename = filename
        # 先通过映射表转换键名，再获取主题
        mapped_theme_key = self.THEME_MAPPING.get(theme_key, theme_key)
        self.theme = THEMES.get(mapped_theme_key, THEMES["商务蓝 (Business Blue)"])
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _set_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["bg_color"]

    def _add_decoration(self, slide):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.4), self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme["accent_color"]
        shape.line.fill.background()

    def add_title_slide(self, title_text, subtitle_text):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_decoration(slide)

        top = Inches(2.5)
        left = Inches(1.0)
        width = Inches(11.3)
        height = Inches(2)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(48)  # 缩小标题字体大小
        p.font.color.rgb = self.theme["title_color"]
        p.font.name = self.theme["font_name"]
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(28)
        p2.font.color.rgb = self.theme["accent_color"]
        p2.font.name = self.theme["font_name"]
        p2.space_before = Pt(20)
        p2.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title_text, points, image_keyword=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_background(slide)
        self._add_decoration(slide)

        # 标题
        txBox_title = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11), Inches(1))
        p_title = txBox_title.text_frame.add_paragraph()
        p_title.text = title_text
        p_title.font.bold = True
        p_title.font.size = Pt(36)
        p_title.font.color.rgb = self.theme["title_color"]
        p_title.font.name = self.theme["font_name"]

        # 装饰线
        shape_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.6), Inches(11), Inches(0.05)
        )
        shape_line.fill.solid()
        shape_line.fill.fore_color.rgb = self.theme["accent_color"]
        shape_line.line.fill.background()

        # 正文
        left_content = Inches(1.0)
        top_content = Inches(2.0)
        width_content = Inches(6.5)
        height_content = Inches(5.0)

        txBox_content = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
        tf_content = txBox_content.text_frame
        tf_content.word_wrap = True

        for point in points:
            p = tf_content.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(22)
            p.font.color.rgb = self.theme["text_color"]
            p.font.name = self.theme["font_name"]
            p.space_after = Pt(20)
            p.line_spacing = 1.3

            # 关键词高亮
            if "核心" in point or "重要" in point or "关键" in point:
                p.font.bold = True
                p.font.color.rgb = self.theme["highlight_color"]

        # 配图
        if image_keyword:
            img_stream = get_image_stream(image_keyword)
            if img_stream:
                img_left = Inches(8.0)
                img_top = Inches(2.0)
                img_height = Inches(3.5)  # 缩小图片高度
                try:
                    pic = slide.shapes.add_picture(img_stream, img_left, img_top, height=img_height)
                    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, img_top, pic.width, pic.height)
                    border.fill.background()
                    border.line.color.rgb = self.theme["accent_color"]
                    border.line.width = Pt(2)
                except Exception as e:
                    print(f"图片插入错误: {e}")

    def save(self):
        self.prs.save(self.filename)


# ===========================
# 4. GUI 界面模块
# ===========================
class PPTApp:
    def __init__(self, root):
        self.root = root
        self.root.title("一键生成 PPT 工具 v1.0")
        self.root.geometry("800x700")

        # --- 顶部：说明与主题选择 ---
        top_frame = ttk.Frame(root, padding=10)
        top_frame.pack(fill=tk.X)

        ttk.Label(top_frame, text="输入内容 (格式：第一行为封面，空行分隔每一页)", font=("Arial", 10, "bold")).pack(
            anchor=tk.W)

        theme_frame = ttk.Frame(top_frame, padding=(0, 5))
        theme_frame.pack(fill=tk.X)
        ttk.Label(theme_frame, text="选择配色主题: ").pack(side=tk.LEFT)

        self.theme_var = tk.StringVar()
        self.theme_combo = ttk.Combobox(theme_frame, textvariable=self.theme_var, state="readonly")
        self.theme_combo['values'] = list(THEMES.keys())
        self.theme_combo.current(1)  # 默认选第二个
        self.theme_combo.pack(side=tk.LEFT, padx=5)

        # --- 中部：文本输入框 ---
        self.text_area = scrolledtext.ScrolledText(root, font=("Arial", 12), height=20)
        self.text_area.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)

        # 预置一些示例文本
        sample_text = """2024年人工智能发展趋势
技术与应用的深度融合

大语言模型的演进
参数规模持续增长
多模态能力增强（文本、图像、视频）
推理成本逐渐降低
开源模型生态日益繁荣

生成式AI的商业落地
核心领域：营销文案与图像生成
辅助编程（Copilot）普及率提升
企业内部知识库问答系统
客户服务自动化

面临的挑战与风险
数据隐私与安全问题
幻觉问题（Hallucination）
版权纠纷
算力资源短缺"""
        self.text_area.insert(tk.END, sample_text)

        # --- 底部：操作区 ---
        bottom_frame = ttk.Frame(root, padding=10)
        bottom_frame.pack(fill=tk.X)

        self.generate_btn = ttk.Button(bottom_frame, text="开始生成 PPT", command=self.on_generate_click)
        self.generate_btn.pack(side=tk.RIGHT, padx=5)

        self.status_label = ttk.Label(bottom_frame, text="准备就绪", foreground="gray")
        self.status_label.pack(side=tk.LEFT)

    def on_generate_click(self):
        # 1. 获取输入
        content = self.text_area.get("1.0", tk.END).strip()
        if not content:
            messagebox.showwarning("提示", "请输入PPT内容！")
            return

        # 2. 选择保存路径
        file_path = filedialog.asksaveasfilename(
            defaultextension=".pptx",
            filetypes=[("PowerPoint Files", "*.pptx")],
            initialfile="my_presentation.pptx"
        )

        if not file_path:
            return

        # 3. 锁定按钮，开始后台线程
        self.generate_btn.config(state=tk.DISABLED)
        self.status_label.config(text="正在分析文本并下载资源，请稍候...", foreground="blue")

        # 开启线程避免界面卡死
        thread = threading.Thread(target=self.run_generation_task, args=(content, file_path))
        thread.start()

    def run_generation_task(self, content, file_path):
        try:
            # A. 分析文本
            analyzer = SimpleTextAnalyzer(content)
            slides_data = analyzer.analyze()

            if not slides_data:
                raise ValueError("无法解析文本内容")

            # B. 初始化生成器
            selected_theme = self.theme_var.get()
            generator = PPTGenerator(file_path, selected_theme)

            total_slides = len(slides_data)

            # C. 逐页生成
            for i, data in enumerate(slides_data):
                # 更新状态文本 (需要在主线程更新UI，这里简化直接调用)
                self.status_label.config(text=f"正在生成第 {i + 1}/{total_slides} 页：{data['title']}...")

                if data["type"] == "title":
                    generator.add_title_slide(data["title"], data["subtitle"])
                elif data["type"] == "content":
                    generator.add_content_slide(data["title"], data["points"], data.get("image_keyword"))

            # D. 保存
            generator.save()

            self.status_label.config(text=f"生成成功！文件已保存至：{file_path}", foreground="green")
            messagebox.showinfo("成功", "PPT 生成完毕！\n您可以直接打开查看。")

        except Exception as e:
            self.status_label.config(text=f"生成失败：{str(e)}", foreground="red")
            messagebox.showerror("错误", f"发生错误：{str(e)}")
        finally:
            self.generate_btn.config(state=tk.NORMAL)


if __name__ == "__main__":
    root = tk.Tk()
    app = PPTApp(root)
    root.mainloop()
