import os
import sys
import time
import traceback
import tempfile

# 添加项目根目录到Python路径
sys.path.insert(0, os.path.abspath('.'))

# 导入转换函数
import ppt2video

# 定义测试PPT文件路径
# 如果没有test_ppt.pptx，可以使用这个脚本生成一个简单的测试PPT
test_ppt_path = "test_ppt.pptx"
output_video_path = "test_output.mp4"

def generate_test_ppt():
    """生成一个简单的测试PPT文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    # 创建演示文稿
    prs = Presentation()
    
    # 添加第一张幻灯片
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = "测试PPT"
    
    # 添加副标题
    subtitle = slide.placeholders[1]
    subtitle.text = "这是一个用于测试PPT转视频功能的演示文稿"
    
    # 添加第二张幻灯片（内容幻灯片）
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = "第一张内容幻灯片"
    
    # 添加内容
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.text = "这是第一张内容幻灯片的文本内容。"
    
    # 添加新段落
    p = text_frame.add_paragraph()
    p.text = "这是第二行文本。"
    p.level = 1
    
    # 添加第三张幻灯片
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "第二张内容幻灯片"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.text = "这是第二张内容幻灯片的文本内容。"
    
    # 保存PPT文件
    prs.save(test_ppt_path)
    print(f"已生成测试PPT文件：{test_ppt_path}")

def main():
    """主测试函数"""
    print("=== PPT转视频调试工具 ===")
    print(f"当前目录：{os.getcwd()}")
    print(f"Python版本：{sys.version}")
    
    # 生成测试PPT
    if not os.path.exists(test_ppt_path):
        print("\n正在生成测试PPT文件...")
        generate_test_ppt()
    else:
        print(f"\n使用已存在的测试PPT文件：{test_ppt_path}")
    
    # 检查测试PPT文件
    if os.path.exists(test_ppt_path):
        print(f"测试PPT文件大小：{os.path.getsize(test_ppt_path)} 字节")
    else:
        print("错误：测试PPT文件不存在！")
        return False
    
    # 清理旧的输出文件
    if os.path.exists(output_video_path):
        os.remove(output_video_path)
        print(f"已清理旧的输出文件：{output_video_path}")
    
    try:
        print("\n=== 开始转换PPT到视频 ===")
        start_time = time.time()
        
        # 调用转换函数
        success = ppt2video.convert_presentation_to_video(
            test_ppt_path, 
            output_video_path, 
            rate=170
        )
        
        end_time = time.time()
        print(f"\n转换完成，耗时：{end_time - start_time:.2f} 秒")
        print(f"转换结果：{'成功' if success else '失败'}")
        
        # 检查输出视频文件
        if os.path.exists(output_video_path):
            file_size = os.path.getsize(output_video_path)
            print(f"输出视频文件：{output_video_path}")
            print(f"视频文件大小：{file_size} 字节")
            print(f"文件是否有效：{'是' if file_size > 10240 else '否（太小）'}")
        else:
            print("错误：输出视频文件不存在！")
        
        return success
        
    except Exception as e:
        print(f"\n转换过程中发生错误：{e}")
        traceback.print_exc()
        return False

if __name__ == "__main__":
    main()